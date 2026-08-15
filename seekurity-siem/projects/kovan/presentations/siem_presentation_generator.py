"""
Seekurity SIEM Presentation Generator
======================================
Generate PowerPoint presentations with SVG diagrams for KOVAN SIEM project.

Features:
- Mermaid to SVG conversion
- SVG component extraction
- PowerPoint slide generation
- Pre-built SIEM architecture diagrams
"""

import json
import os
import re
import subprocess
import tempfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Dict, List, Optional, Tuple
from xml.etree import ElementTree as ET

# Check for python-pptx
try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# =============================================================================
# MERMAID DIAGRAMS FOR KOVAN SIEM
# =============================================================================

SIEM_DIAGRAMS = {
    "architecture_overview": """
%%{init: {'theme': 'base', 'themeVariables': { 'fontSize': '14px' }, 'flowchart': { 'useMaxWidth': true }}}%%
flowchart TB
    subgraph Sources["Log Sources"]
        style Sources fill:#e3f2fd,stroke:#1976d2
        FW[Firewalls]
        VPN[VPN Gateways]
        NAC[NAC Systems]
        IPS[IPS/IDS]
    end

    subgraph Collection["Log Collection Layer"]
        style Collection fill:#fff3e0,stroke:#f57c00
        Syslog[Syslog Collector]
        Agent[Log Agents]
        API[API Collectors]
    end

    subgraph Processing["Processing Engine"]
        style Processing fill:#e8f5e9,stroke:#43a047
        Parser[Log Parser]
        Norm[Normalizer]
        Enrich[Enrichment]
    end

    subgraph Detection["Detection Engine"]
        style Detection fill:#fce4ec,stroke:#e91e63
        Rules[Detection Rules]
        Corr[Correlation]
        ML[ML Analytics]
    end

    subgraph Response["Response & Storage"]
        style Response fill:#f3e5f5,stroke:#9c27b0
        Alert[Alert Manager]
        SOAR[SOAR Integration]
        Store[(Elasticsearch)]
    end

    Sources --> Collection
    Collection --> Processing
    Processing --> Detection
    Detection --> Response
""",
    "log_flow": """
%%{init: {'theme': 'base', 'themeVariables': { 'fontSize': '14px' }, 'flowchart': { 'useMaxWidth': true }}}%%
flowchart LR
    subgraph Input["Input"]
        style Input fill:#e3f2fd,stroke:#1976d2
        Raw[Raw Logs]
    end

    subgraph Parse["Parsing"]
        style Parse fill:#fff3e0,stroke:#f57c00
        Regex[Regex Engine]
        Fields[Field Extraction]
    end

    subgraph Detect["Detection"]
        style Detect fill:#fce4ec,stroke:#e91e63
        Match[Rule Matching]
        Thresh[Threshold Check]
    end

    subgraph Output["Output"]
        style Output fill:#e8f5e9,stroke:#43a047
        Alert[Alerts]
        Store[Storage]
    end

    Raw --> Regex --> Fields --> Match --> Thresh --> Alert
    Thresh --> Store
""",
    "detection_rules": """
%%{init: {'theme': 'base', 'themeVariables': { 'fontSize': '14px' }, 'flowchart': { 'useMaxWidth': true }}}%%
flowchart TB
    subgraph Rules["Detection Rules (28)"]
        style Rules fill:#e3f2fd,stroke:#1976d2

        subgraph FW["Firewall (6)"]
            style FW fill:#ffcdd2,stroke:#e53935
            FW1[Port Scan]
            FW2[C2 Traffic]
            FW3[DB Access]
        end

        subgraph VPN["VPN (4)"]
            style VPN fill:#c8e6c9,stroke:#43a047
            VPN1[Brute Force]
            VPN2[Auth Failure]
        end

        subgraph NAC["NAC (4)"]
            style NAC fill:#fff9c4,stroke:#f9a825
            NAC1[Rogue Device]
            NAC2[MAC Spoof]
        end

        subgraph WAF["WAF/IPS (6)"]
            style WAF fill:#d1c4e9,stroke:#7b1fa2
            WAF1[SQLi]
            WAF2[XSS]
            WAF3[RCE]
        end
    end

    subgraph Severity["Severity Levels"]
        style Severity fill:#fff3e0,stroke:#f57c00
        C[CRITICAL: 5]
        H[HIGH: 11]
        M[MEDIUM: 9]
        L[LOW: 3]
    end

    Rules --> Severity
""",
    "kovan_network": """
%%{init: {'theme': 'base', 'themeVariables': { 'fontSize': '14px' }, 'flowchart': { 'useMaxWidth': true }}}%%
flowchart TB
    subgraph Internet["Internet Zone"]
        style Internet fill:#ffcdd2,stroke:#e53935
        Ext[External Users]
    end

    subgraph DMZ["DMZ"]
        style DMZ fill:#fff9c4,stroke:#f9a825
        WAF[WAPPLES WAF]
        Web[Web Servers]
    end

    subgraph Security["Security Layer"]
        style Security fill:#c8e6c9,stroke:#43a047
        FW1[SECUI MF2]
        IPS1[SECUI IPS]
        DDoS[WINS SNIPER]
    end

    subgraph Internal["Internal Network"]
        style Internal fill:#e3f2fd,stroke:#1976d2
        NAC[Genians NAC]
        VPN[VPN Gateway]
        DB[(Databases)]
    end

    subgraph SIEM["Seekurity SIEM"]
        style SIEM fill:#f3e5f5,stroke:#9c27b0
        Collector[Log Collector]
        Engine[Detection Engine]
        Dashboard[Dashboard]
    end

    Ext --> WAF --> FW1 --> Internal
    FW1 --> IPS1
    DDoS --> FW1
    Security --> Collector
    Internal --> Collector
    Collector --> Engine --> Dashboard
""",
    "vendor_coverage": """
%%{init: {'theme': 'base', 'themeVariables': { 'fontSize': '14px' }, 'flowchart': { 'useMaxWidth': true }}}%%
flowchart TB
    subgraph Verified["Verified Parsers (8)"]
        style Verified fill:#c8e6c9,stroke:#43a047
        J[Juniper SSG]
        F[FortiGate]
        S[SECUI MF2]
        G[Genians NAC]
    end

    subgraph Unverified["Pending Verification (14)"]
        style Unverified fill:#fff9c4,stroke:#f9a825
        W[WINS SNIPER]
        P[Penta WAPPLES]
        N[NexG VForce]
        H[Hunesion]
    end

    subgraph Sources["Total Log Sources"]
        style Sources fill:#e3f2fd,stroke:#1976d2
        T1[Firewalls: 20+]
        T2[VPN: 15+]
        T3[IPS/DDoS: 6]
        T4[NAC/Other: 10+]
    end

    Verified --> Sources
    Unverified --> Sources
""",
    "alert_workflow": """
%%{init: {'theme': 'base', 'themeVariables': { 'fontSize': '14px' }, 'flowchart': { 'useMaxWidth': true }}}%%
flowchart LR
    subgraph Trigger["Event Trigger"]
        style Trigger fill:#e3f2fd,stroke:#1976d2
        Log[Security Log]
    end

    subgraph Analysis["Analysis"]
        style Analysis fill:#fff3e0,stroke:#f57c00
        Parse[Parse & Normalize]
        Rule[Rule Matching]
        MITRE[MITRE Mapping]
    end

    subgraph Alert["Alert Generation"]
        style Alert fill:#fce4ec,stroke:#e91e63
        Sev{Severity?}
        Crit[CRITICAL]
        High[HIGH]
        Med[MEDIUM]
    end

    subgraph Response["Response"]
        style Response fill:#c8e6c9,stroke:#43a047
        Auto[Auto-Response]
        Notify[Notification]
        Ticket[Ticket Creation]
    end

    Log --> Parse --> Rule --> MITRE --> Sev
    Sev -->|Critical| Crit --> Auto
    Sev -->|High| High --> Notify
    Sev -->|Medium| Med --> Ticket
""",
}


# =============================================================================
# SVG GENERATOR
# =============================================================================


@dataclass
class SVGComponent:
    """Represents an SVG component"""

    id: str
    tag: str
    attributes: Dict[str, str]
    children: List["SVGComponent"] = field(default_factory=list)
    text: Optional[str] = None


class MermaidToSVG:
    """Convert Mermaid diagrams to SVG"""

    def __init__(self):
        self.mmdc_path = self._find_mmdc()

    def _find_mmdc(self) -> Optional[str]:
        """Find mermaid-cli (mmdc) executable"""
        # Try common locations
        possible_paths = [
            "mmdc",
            "npx mmdc",
            r"C:\Users\{}\AppData\Roaming\npm\mmdc.cmd".format(
                os.environ.get("USERNAME", "")
            ),
            "/usr/local/bin/mmdc",
            "/usr/bin/mmdc",
        ]

        for path in possible_paths:
            try:
                result = subprocess.run(
                    f"{path} --version", shell=True, capture_output=True, timeout=10
                )
                if result.returncode == 0:
                    return path
            except:
                continue
        return None

    def convert(self, mermaid_code: str, output_path: str) -> bool:
        """Convert Mermaid code to SVG file"""
        if not self.mmdc_path:
            print(
                "Warning: mermaid-cli not found. Install with: npm install -g @mermaid-js/mermaid-cli"
            )
            return self._create_placeholder_svg(mermaid_code, output_path)

        # Write mermaid code to temp file
        with tempfile.NamedTemporaryFile(mode="w", suffix=".mmd", delete=False) as f:
            f.write(mermaid_code)
            temp_input = f.name

        try:
            cmd = (
                f'{self.mmdc_path} -i "{temp_input}" -o "{output_path}" -b transparent'
            )
            result = subprocess.run(cmd, shell=True, capture_output=True, timeout=30)
            return result.returncode == 0
        except Exception as e:
            print(f"Error converting Mermaid: {e}")
            return self._create_placeholder_svg(mermaid_code, output_path)
        finally:
            os.unlink(temp_input)

    def _create_placeholder_svg(self, mermaid_code: str, output_path: str) -> bool:
        """Create a placeholder SVG with text representation"""
        # Extract title from mermaid code
        title = "SIEM Diagram"
        if "subgraph" in mermaid_code:
            match = re.search(r'subgraph\s+(\w+)\["([^"]+)"\]', mermaid_code)
            if match:
                title = match.group(2)

        svg_content = f"""<?xml version="1.0" encoding="UTF-8"?>
<svg xmlns="http://www.w3.org/2000/svg" width="800" height="600" viewBox="0 0 800 600">
  <defs>
    <linearGradient id="bg" x1="0%" y1="0%" x2="0%" y2="100%">
      <stop offset="0%" style="stop-color:#f5f5f5"/>
      <stop offset="100%" style="stop-color:#e0e0e0"/>
    </linearGradient>
  </defs>

  <!-- Background -->
  <rect width="800" height="600" fill="url(#bg)" rx="10"/>

  <!-- Border -->
  <rect x="5" y="5" width="790" height="590" fill="none" stroke="#1976d2" stroke-width="2" rx="8"/>

  <!-- Title -->
  <text x="400" y="50" text-anchor="middle" font-family="Arial, sans-serif" font-size="24" font-weight="bold" fill="#1976d2">
    {title}
  </text>

  <!-- Seekurity Logo placeholder -->
  <rect x="300" y="100" width="200" height="80" fill="#1976d2" rx="5"/>
  <text x="400" y="150" text-anchor="middle" font-family="Arial, sans-serif" font-size="20" fill="white">
    Seekurity SIEM
  </text>

  <!-- Diagram placeholder boxes -->
  <g id="diagram_components">
    <rect x="100" y="220" width="150" height="60" fill="#e3f2fd" stroke="#1976d2" rx="5"/>
    <text x="175" y="255" text-anchor="middle" font-family="Arial, sans-serif" font-size="14">Log Sources</text>

    <rect x="325" y="220" width="150" height="60" fill="#fff3e0" stroke="#f57c00" rx="5"/>
    <text x="400" y="255" text-anchor="middle" font-family="Arial, sans-serif" font-size="14">Processing</text>

    <rect x="550" y="220" width="150" height="60" fill="#fce4ec" stroke="#e91e63" rx="5"/>
    <text x="625" y="255" text-anchor="middle" font-family="Arial, sans-serif" font-size="14">Detection</text>

    <!-- Arrows -->
    <path d="M 250 250 L 320 250" stroke="#666" stroke-width="2" marker-end="url(#arrow)"/>
    <path d="M 475 250 L 545 250" stroke="#666" stroke-width="2" marker-end="url(#arrow)"/>
  </g>

  <!-- Arrow marker -->
  <defs>
    <marker id="arrow" markerWidth="10" markerHeight="10" refX="9" refY="3" orient="auto">
      <path d="M0,0 L0,6 L9,3 z" fill="#666"/>
    </marker>
  </defs>

  <!-- Footer -->
  <text x="400" y="550" text-anchor="middle" font-family="Arial, sans-serif" font-size="12" fill="#666">
    KOVAN Project - Security Information and Event Management
  </text>
</svg>"""

        with open(output_path, "w", encoding="utf-8") as f:
            f.write(svg_content)
        return True


# =============================================================================
# SVG COMPONENT EXTRACTOR
# =============================================================================


class SVGExtractor:
    """Extract and manipulate SVG components"""

    SVG_NS = {"svg": "http://www.w3.org/2000/svg"}

    def __init__(self, svg_path: str):
        self.tree = ET.parse(svg_path)
        self.root = self.tree.getroot()

    def list_components(self) -> List[Dict]:
        """List all components in SVG"""
        components = []
        for elem in self.root.iter():
            tag = elem.tag.replace("{http://www.w3.org/2000/svg}", "")
            elem_id = elem.get("id", "")
            if elem_id or tag in ["g", "rect", "circle", "path", "text"]:
                components.append(
                    {
                        "tag": tag,
                        "id": elem_id,
                        "class": elem.get("class", ""),
                        "attributes": dict(elem.attrib),
                    }
                )
        return components

    def extract_component(self, component_id: str) -> Optional[str]:
        """Extract a specific component by ID"""
        for elem in self.root.iter():
            if elem.get("id") == component_id:
                return ET.tostring(elem, encoding="unicode")
        return None

    def get_groups(self) -> List[Dict]:
        """Get all group elements (layers)"""
        groups = []
        for elem in self.root.iter():
            tag = elem.tag.replace("{http://www.w3.org/2000/svg}", "")
            if tag == "g":
                groups.append(
                    {
                        "id": elem.get("id", "unnamed"),
                        "class": elem.get("class", ""),
                        "children": len(list(elem)),
                    }
                )
        return groups


# =============================================================================
# POWERPOINT GENERATOR
# =============================================================================


class SIEMPresentationGenerator:
    """Generate PowerPoint presentations for Seekurity SIEM"""

    # Color scheme
    COLORS = {
        "primary": RGBColor(25, 118, 210) if PPTX_AVAILABLE else None,  # Blue
        "secondary": RGBColor(245, 124, 0) if PPTX_AVAILABLE else None,  # Orange
        "success": RGBColor(67, 160, 71) if PPTX_AVAILABLE else None,  # Green
        "danger": RGBColor(229, 57, 53) if PPTX_AVAILABLE else None,  # Red
        "warning": RGBColor(251, 192, 45) if PPTX_AVAILABLE else None,  # Yellow
        "dark": RGBColor(33, 33, 33) if PPTX_AVAILABLE else None,  # Dark gray
        "light": RGBColor(245, 245, 245) if PPTX_AVAILABLE else None,  # Light gray
    }

    def __init__(self, output_path: str):
        if not PPTX_AVAILABLE:
            raise ImportError(
                "python-pptx is required. Install with: pip install python-pptx"
            )

        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        self.prs.slide_height = Inches(7.5)
        self.output_path = output_path
        self.mermaid_converter = MermaidToSVG()

    def add_title_slide(self, title: str, subtitle: str):
        """Add title slide"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # Background
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.COLORS["primary"]
        background.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(48)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        title_para.alignment = PP_ALIGN.CENTER

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(12.333), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = subtitle
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = RGBColor(255, 255, 255)
        subtitle_para.alignment = PP_ALIGN.CENTER

    def add_content_slide(
        self, title: str, content: List[str], diagram_key: Optional[str] = None
    ):
        """Add content slide with optional diagram"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # Title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.COLORS["primary"]
        title_bar.line.fill.background()

        # Title text
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        # Content area
        if diagram_key and diagram_key in SIEM_DIAGRAMS:
            # Left side: bullet points
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.3), Inches(5), Inches(5.5)
            )
            self._add_bullet_points(content_box.text_frame, content)

            # Right side: diagram placeholder
            diagram_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(6),
                Inches(1.3),
                Inches(6.833),
                Inches(5.5),
            )
            diagram_box.fill.solid()
            diagram_box.fill.fore_color.rgb = self.COLORS["light"]

            # Add diagram label
            label_box = slide.shapes.add_textbox(
                Inches(6), Inches(3.5), Inches(6.833), Inches(0.5)
            )
            label_frame = label_box.text_frame
            label_para = label_frame.paragraphs[0]
            label_para.text = f"[Diagram: {diagram_key}]"
            label_para.font.size = Pt(14)
            label_para.font.color.rgb = self.COLORS["dark"]
            label_para.alignment = PP_ALIGN.CENTER
        else:
            # Full width content
            content_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.3), Inches(12.333), Inches(5.5)
            )
            self._add_bullet_points(content_box.text_frame, content)

    def _add_bullet_points(self, text_frame, points: List[str]):
        """Add bullet points to text frame"""
        text_frame.word_wrap = True

        for i, point in enumerate(points):
            if i == 0:
                para = text_frame.paragraphs[0]
            else:
                para = text_frame.add_paragraph()

            para.text = f"• {point}"
            para.font.size = Pt(18)
            para.font.color.rgb = self.COLORS["dark"]
            para.space_after = Pt(12)

    def add_table_slide(self, title: str, headers: List[str], rows: List[List[str]]):
        """Add slide with table"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.COLORS["primary"]
        title_bar.line.fill.background()

        # Title text
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        # Table
        num_rows = len(rows) + 1  # +1 for header
        num_cols = len(headers)

        table = slide.shapes.add_table(
            num_rows,
            num_cols,
            Inches(0.5),
            Inches(1.3),
            Inches(12.333),
            Inches(0.5 * num_rows),
        ).table

        # Header row
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.COLORS["primary"]
            para = cell.text_frame.paragraphs[0]
            para.font.bold = True
            para.font.size = Pt(14)
            para.font.color.rgb = RGBColor(255, 255, 255)

        # Data rows
        for row_idx, row in enumerate(rows):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = value
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(12)

    def add_statistics_slide(self, title: str, stats: Dict[str, any]):
        """Add slide with statistics boxes"""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)

        # Title bar
        title_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1)
        )
        title_bar.fill.solid()
        title_bar.fill.fore_color.rgb = self.COLORS["primary"]
        title_bar.line.fill.background()

        # Title text
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.2), Inches(12.333), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = title
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        # Statistics boxes
        colors = [
            self.COLORS["primary"],
            self.COLORS["success"],
            self.COLORS["secondary"],
            self.COLORS["danger"],
        ]

        items = list(stats.items())[:4]  # Max 4 stats
        box_width = Inches(2.8)
        start_x = Inches(0.5)
        gap = Inches(0.4)

        for i, (label, value) in enumerate(items):
            x = start_x + (box_width + gap) * i

            # Box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2), box_width, Inches(2.5)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = colors[i % len(colors)]

            # Value
            value_box = slide.shapes.add_textbox(x, Inches(2.3), box_width, Inches(1))
            value_frame = value_box.text_frame
            value_para = value_frame.paragraphs[0]
            value_para.text = str(value)
            value_para.font.size = Pt(48)
            value_para.font.bold = True
            value_para.font.color.rgb = RGBColor(255, 255, 255)
            value_para.alignment = PP_ALIGN.CENTER

            # Label
            label_box = slide.shapes.add_textbox(x, Inches(3.5), box_width, Inches(0.8))
            label_frame = label_box.text_frame
            label_para = label_frame.paragraphs[0]
            label_para.text = label
            label_para.font.size = Pt(16)
            label_para.font.color.rgb = RGBColor(255, 255, 255)
            label_para.alignment = PP_ALIGN.CENTER

    def save(self):
        """Save presentation"""
        self.prs.save(self.output_path)
        print(f"Presentation saved to: {self.output_path}")


# =============================================================================
# MAIN FUNCTIONS
# =============================================================================


def generate_siem_presentation(output_path: str):
    """Generate complete SIEM presentation"""
    if not PPTX_AVAILABLE:
        print("Error: python-pptx not installed. Run: pip install python-pptx")
        return

    gen = SIEMPresentationGenerator(output_path)

    # Slide 1: Title
    gen.add_title_slide(
        "Seekurity SIEM",
        "KOVAN Security Information and Event Management\n보안 정보 및 이벤트 관리 시스템",
    )

    # Slide 2: Architecture Overview
    gen.add_content_slide(
        "System Architecture",
        [
            "Multi-tier log collection architecture",
            "Real-time log parsing and normalization",
            "Advanced detection rule engine",
            "MITRE ATT&CK framework mapping",
            "Automated response and alerting",
        ],
        diagram_key="architecture_overview",
    )

    # Slide 3: Log Sources
    gen.add_table_slide(
        "Supported Log Sources",
        ["Category", "Vendor", "Models", "Status"],
        [
            ["Firewall", "SECUI", "MF2 1510/1520, BLUEMAX NGF", "✓ Verified"],
            ["Firewall", "Juniper", "SSG-550M, SSG-320M", "✓ Verified"],
            ["VPN", "FortiGate", "100D", "✓ Verified"],
            ["NAC", "Genians", "S10_R1, C10_R1", "✓ Verified"],
            ["IPS/DDoS", "WINS", "SNIPER DDX/ONE-d", "Pending"],
            ["WAF", "Penta", "WAPPLES-2700", "Pending"],
        ],
    )

    # Slide 4: Detection Rules Statistics
    gen.add_statistics_slide(
        "Detection Rules Overview",
        {"Total Rules": 28, "CRITICAL": 5, "HIGH": 11, "Categories": 8},
    )

    # Slide 5: Detection Categories
    gen.add_content_slide(
        "Detection Rule Categories",
        [
            "Firewall: Port scans, C2 traffic, unauthorized DB access",
            "VPN: Brute force attacks, authentication failures",
            "NAC: Rogue devices, MAC spoofing, policy violations",
            "IPS/IDS: SQL injection, XSS, exploit signatures",
            "DDoS: Flood attacks, traffic anomalies",
            "WAF: Command injection, directory traversal",
            "Authentication: Failed logins, privilege escalation",
            "Exfiltration: Large data transfers, DNS tunneling",
        ],
        diagram_key="detection_rules",
    )

    # Slide 6: MITRE ATT&CK Mapping
    gen.add_table_slide(
        "MITRE ATT&CK Coverage",
        ["Tactic", "Technique ID", "Detection Rules"],
        [
            ["Initial Access", "T1190, T1133", "FW-005, FW-006, VPN-003"],
            ["Credential Access", "T1110", "VPN-001, VPN-002, AUTH-001"],
            ["Defense Evasion", "T1036, T1562", "NAC-002, NAC-004"],
            ["Command & Control", "T1571, T1071", "FW-002, EXFIL-002"],
            ["Exfiltration", "T1048", "FW-004, EXFIL-001"],
            ["Impact", "T1498", "DDOS-001, DDOS-002"],
        ],
    )

    # Slide 7: Alert Workflow
    gen.add_content_slide(
        "Alert Processing Workflow",
        [
            "Log ingestion from 50+ sources",
            "Real-time parsing with vendor-specific regex",
            "Rule matching against 28 detection rules",
            "Automatic MITRE ATT&CK classification",
            "Severity-based response actions",
            "Integration with ticketing and SOAR systems",
        ],
        diagram_key="alert_workflow",
    )

    # Slide 8: Network Coverage
    gen.add_content_slide(
        "KOVAN Network Coverage",
        [
            "Internet Zone: External user access monitoring",
            "DMZ: Web application firewall logs",
            "Security Layer: Firewall, IPS, DDoS protection",
            "Internal Network: NAC, VPN, database access",
            "Full visibility across all network segments",
        ],
        diagram_key="kovan_network",
    )

    gen.save()


def export_diagrams(output_dir: str):
    """Export all Mermaid diagrams to SVG files"""
    converter = MermaidToSVG()
    Path(output_dir).mkdir(parents=True, exist_ok=True)

    for name, mermaid_code in SIEM_DIAGRAMS.items():
        output_path = os.path.join(output_dir, f"{name}.svg")
        success = converter.convert(mermaid_code, output_path)
        status = "✓" if success else "✗"
        print(f"{status} {name}.svg")


def list_diagram_components(svg_path: str):
    """List components in an SVG file"""
    extractor = SVGExtractor(svg_path)

    print(f"\nComponents in {svg_path}:")
    print("-" * 50)

    groups = extractor.get_groups()
    print(f"\nGroups (Layers): {len(groups)}")
    for g in groups:
        print(f"  - {g['id']} ({g['children']} children)")

    print("\nAll Components:")
    for comp in extractor.list_components()[:20]:  # First 20
        print(f"  {comp['tag']:10} id={comp['id']}")


# =============================================================================
# CLI
# =============================================================================

if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Seekurity SIEM Presentation Generator")
        print("=" * 40)
        print("\nUsage:")
        print("  python siem_presentation_generator.py presentation <output.pptx>")
        print("  python siem_presentation_generator.py diagrams <output_dir>")
        print("  python siem_presentation_generator.py components <svg_file>")
        print("\nExamples:")
        print("  python siem_presentation_generator.py presentation ./KOVAN_SIEM.pptx")
        print("  python siem_presentation_generator.py diagrams ./diagrams")
        sys.exit(0)

    command = sys.argv[1]

    if command == "presentation":
        output = sys.argv[2] if len(sys.argv) > 2 else "KOVAN_SIEM_Presentation.pptx"
        generate_siem_presentation(output)

    elif command == "diagrams":
        output_dir = sys.argv[2] if len(sys.argv) > 2 else "./diagrams"
        export_diagrams(output_dir)

    elif command == "components":
        if len(sys.argv) < 3:
            print("Error: SVG file path required")
            sys.exit(1)
        list_diagram_components(sys.argv[2])

    else:
        print(f"Unknown command: {command}")
        sys.exit(1)
