"""
Seekurity SIEM Professional Presentation Generator
Generates enterprise-grade PowerPoint presentations following IBM/HP/Microsoft styling.
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# =============================================================================
# BRAND COLORS (IBM-inspired)
# =============================================================================
COLORS = {
    "primary_blue": RGBColor(0x0F, 0x62, 0xFE),  # #0f62fe
    "secondary_blue": RGBColor(0x03, 0x53, 0xE9),  # #0353e9
    "dark_gray": RGBColor(0x16, 0x16, 0x16),  # #161616
    "medium_gray": RGBColor(0x52, 0x52, 0x52),  # #525252
    "light_gray": RGBColor(0x6F, 0x6F, 0x6F),  # #6f6f6f
    "background": RGBColor(0xFF, 0xFF, 0xFF),  # #ffffff
    "alt_background": RGBColor(0xF4, 0xF4, 0xF4),  # #f4f4f4
    "success_green": RGBColor(0x19, 0x80, 0x38),  # #198038
    "warning_yellow": RGBColor(0xF1, 0xC2, 0x1B),  # #f1c21b
    "error_red": RGBColor(0xDA, 0x1E, 0x28),  # #da1e28
    "white": RGBColor(0xFF, 0xFF, 0xFF),  # #ffffff
    "table_header": RGBColor(0x0F, 0x62, 0xFE),  # #0f62fe
    "table_row_alt": RGBColor(0xE8, 0xE8, 0xE8),  # #e8e8e8
}

# =============================================================================
# SLIDE DIMENSIONS (16:9 Widescreen)
# =============================================================================
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# =============================================================================
# HELPER FUNCTIONS
# =============================================================================


def set_shape_fill(shape, color):
    """Set solid fill color for a shape."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_shape_transparency(shape, transparency):
    """Set transparency for a shape (0-100)."""
    fill = shape.fill
    fill.solid()
    # Transparency is set via alpha channel
    alpha = int((100 - transparency) * 1000)  # Convert to EMUs


def add_text_frame(
    shape, text, font_size=18, font_bold=False, font_color=None, alignment=PP_ALIGN.LEFT
):
    """Add formatted text to a shape."""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = font_bold
    run.font.name = "Segoe UI"
    if font_color:
        run.font.color.rgb = font_color
    return tf


def add_paragraph(
    text_frame,
    text,
    font_size=18,
    font_bold=False,
    font_color=None,
    alignment=PP_ALIGN.LEFT,
    space_before=0,
    space_after=0,
    level=0,
):
    """Add a new paragraph to a text frame."""
    p = text_frame.add_paragraph()
    p.alignment = alignment
    p.level = level
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = font_bold
    run.font.name = "Segoe UI"
    if font_color:
        run.font.color.rgb = font_color
    return p


def create_title_slide(prs, title, subtitle, customer_name, date, company="Seekurity"):
    """Create a title slide with blue background."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Blue background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    set_shape_fill(background, COLORS["primary_blue"])
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(2.2), Inches(11.8), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = COLORS["white"]
    run.font.name = "Segoe UI"

    # Subtitle
    add_paragraph(
        tf,
        subtitle,
        font_size=28,
        font_color=COLORS["white"],
        alignment=PP_ALIGN.CENTER,
        space_before=12,
    )

    # Customer name
    customer_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.2), Inches(11.8), Inches(0.6)
    )
    tf2 = customer_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = customer_name
    run2.font.size = Pt(24)
    run2.font.bold = True
    run2.font.color.rgb = COLORS["white"]
    run2.font.name = "Segoe UI"

    # Company and date
    footer_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(5.5), Inches(11.8), Inches(0.8)
    )
    tf3 = footer_box.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = f"Prepared by {company}"
    run3.font.size = Pt(16)
    run3.font.color.rgb = COLORS["white"]
    run3.font.name = "Segoe UI"
    add_paragraph(
        tf3,
        date,
        font_size=14,
        font_color=COLORS["white"],
        alignment=PP_ALIGN.CENTER,
        space_before=6,
    )

    return slide


def create_section_divider(prs, section_number, section_title):
    """Create a section divider slide with gray background."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Gray background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    set_shape_fill(background, COLORS["alt_background"])
    background.line.fill.background()

    # Section number
    num_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(2.5), Inches(11.8), Inches(1)
    )
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{section_number:02d}"
    run.font.size = Pt(72)
    run.font.bold = True
    run.font.color.rgb = COLORS["primary_blue"]
    run.font.name = "Segoe UI"

    # Section title
    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(3.8), Inches(11.8), Inches(0.8)
    )
    tf2 = title_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = section_title
    run2.font.size = Pt(36)
    run2.font.bold = False
    run2.font.color.rgb = COLORS["dark_gray"]
    run2.font.name = "Segoe UI"

    return slide


def create_content_slide(prs, title, subtitle=None):
    """Create a content slide with title and optional subtitle."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(0.5), Inches(11.8), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = COLORS["primary_blue"]
    run.font.name = "Segoe UI"

    # Subtitle
    if subtitle:
        add_paragraph(
            tf, subtitle, font_size=20, font_color=COLORS["medium_gray"], space_before=6
        )

    # Bottom line accent
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.35), Inches(1.5), Inches(0.04)
    )
    set_shape_fill(line, COLORS["primary_blue"])
    line.line.fill.background()

    return slide


def add_table(
    slide,
    data,
    left=0.75,
    top=1.6,
    width=11.8,
    row_height=0.45,
    header_color=None,
    font_size=12,
):
    """Add a formatted table to a slide."""
    rows = len(data)
    cols = len(data[0]) if data else 0

    if rows == 0 or cols == 0:
        return None

    # Calculate column width
    col_width = width / cols
    table_height = rows * row_height

    table = slide.shapes.add_table(
        rows, cols, Inches(left), Inches(top), Inches(width), Inches(table_height)
    ).table

    # Style table
    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_text)

            # Format cell text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Segoe UI"

                if i == 0:  # Header row
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = COLORS["white"]
                else:
                    paragraph.font.color.rgb = COLORS["dark_gray"]

            # Cell background
            if i == 0:  # Header
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color or COLORS["table_header"]
            elif i % 2 == 0:  # Alternating rows
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["table_row_alt"]
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["white"]

            # Vertical alignment
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table


def add_bullet_points(
    slide,
    points,
    left=0.75,
    top=1.6,
    width=11.8,
    height=5,
    font_size=18,
    bullet_color=None,
):
    """Add bullet points to a slide."""
    textbox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = textbox.text_frame
    tf.word_wrap = True

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.level = point.get("level", 0)
        p.space_before = Pt(6)
        p.space_after = Pt(6)

        run = p.add_run()
        run.text = (
            f"• {point['text']}"
            if p.level == 0
            else f"  {'  ' * p.level}– {point['text']}"
        )
        run.font.size = Pt(font_size - (p.level * 2))
        run.font.name = "Segoe UI"
        run.font.color.rgb = point.get("color", COLORS["dark_gray"])
        run.font.bold = point.get("bold", False)

    return textbox


def create_thank_you_slide(
    prs,
    company="Seekurity",
    tagline="Enterprise Security Solutions",
    email="contact@seekurity.com",
    phone="+1 (555) 123-4567",
    website="www.seekurity.com",
):
    """Create a thank you/closing slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Blue background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    set_shape_fill(background, COLORS["primary_blue"])
    background.line.fill.background()

    # Thank you text
    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(2.2), Inches(11.8), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Thank You"
    run.font.size = Pt(56)
    run.font.bold = True
    run.font.color.rgb = COLORS["white"]
    run.font.name = "Segoe UI"

    # Company name
    add_paragraph(
        tf,
        company,
        font_size=28,
        font_bold=True,
        font_color=COLORS["white"],
        alignment=PP_ALIGN.CENTER,
        space_before=24,
    )
    add_paragraph(
        tf,
        tagline,
        font_size=18,
        font_color=COLORS["white"],
        alignment=PP_ALIGN.CENTER,
        space_before=6,
    )

    # Contact info
    contact_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(5.2), Inches(11.8), Inches(1.2)
    )
    tf2 = contact_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = email
    run2.font.size = Pt(14)
    run2.font.color.rgb = COLORS["white"]
    run2.font.name = "Segoe UI"

    add_paragraph(
        tf2,
        phone,
        font_size=14,
        font_color=COLORS["white"],
        alignment=PP_ALIGN.CENTER,
        space_before=4,
    )
    add_paragraph(
        tf2,
        website,
        font_size=14,
        font_color=COLORS["white"],
        alignment=PP_ALIGN.CENTER,
        space_before=4,
    )

    return slide


# =============================================================================
# KICKOFF PRESENTATION GENERATOR
# =============================================================================


def generate_kickoff_presentation(
    output_path, customer_name="[CUSTOMER_NAME]", date="[DATE]"
):
    """Generate the Project Kickoff presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # =========================================================================
    # TITLE SLIDE
    # =========================================================================
    create_title_slide(
        prs, "SIEM Implementation", "Project Kickoff", customer_name, date
    )

    # =========================================================================
    # AGENDA
    # =========================================================================
    slide = create_content_slide(prs, "Agenda")
    agenda_data = [
        ["#", "Topic", "Duration"],
        ["1", "Introduction & Objectives", "10 min"],
        ["2", "Project Scope & Deliverables", "15 min"],
        ["3", "Technical Architecture", "20 min"],
        ["4", "Implementation Approach", "15 min"],
        ["5", "Timeline & Milestones", "10 min"],
        ["6", "Governance & Communication", "10 min"],
        ["7", "Risks & Mitigation", "10 min"],
        ["8", "Next Steps & Q&A", "10 min"],
    ]
    add_table(slide, agenda_data)

    # Add total duration note
    note_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(5.8), Inches(11.8), Inches(0.4)
    )
    tf = note_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Total Duration: ~90 minutes"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    # =========================================================================
    # SECTION 1: INTRODUCTION & OBJECTIVES
    # =========================================================================
    create_section_divider(prs, 1, "Introduction & Objectives")

    # Meeting Objectives slide
    slide = create_content_slide(prs, "Meeting Objectives")
    objectives = [
        {"text": "Aligned on project goals and success criteria", "bold": True},
        {"text": "Confirmed scope, deliverables, and exclusions", "bold": True},
        {
            "text": "Reviewed technical architecture and integration points",
            "bold": True,
        },
        {"text": "Agreed on timeline, milestones, and dependencies", "bold": True},
        {"text": "Established governance model and escalation paths", "bold": True},
        {"text": "Identified risks and mitigation strategies", "bold": True},
    ]
    add_bullet_points(slide, objectives)

    # Project Team slide
    slide = create_content_slide(prs, "Project Team", "Seekurity Team")
    seekurity_team = [
        ["Role", "Name", "Responsibility"],
        ["Project Manager", "[NAME]", "Overall delivery, stakeholder management"],
        ["Lead Engineer", "[NAME]", "Technical architecture, implementation"],
        ["Security Analyst", "[NAME]", "Use case development, tuning"],
        ["Support Lead", "[NAME]", "Knowledge transfer, documentation"],
    ]
    add_table(slide, seekurity_team, top=1.6, row_height=0.5)

    # Customer Team
    customer_label = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.0), Inches(11.8), Inches(0.4)
    )
    tf = customer_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"{customer_name} Team"
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    customer_team = [
        ["Role", "Name", "Responsibility"],
        ["Project Sponsor", "[NAME]", "Executive oversight, decisions"],
        ["Technical Lead", "[NAME]", "Infrastructure, access, coordination"],
        ["Security Lead", "[NAME]", "Requirements, validation, FVT"],
    ]
    add_table(slide, customer_team, top=4.5, row_height=0.5)

    # =========================================================================
    # SECTION 2: PROJECT SCOPE & DELIVERABLES
    # =========================================================================
    create_section_divider(prs, 2, "Project Scope & Deliverables")

    # Scope Overview slide
    slide = create_content_slide(prs, "Scope Overview", "In Scope")
    scope_data = [
        ["Category", "Details"],
        ["Platform", "[SIEM Platform] deployment and configuration"],
        ["Log Sources", "[X] sources across [Y] categories"],
        ["Use Cases", "[Z] detection rules and alerts"],
        ["Integrations", "Ticketing, SOAR, threat intelligence"],
        ["Documentation", "Runbooks, architecture diagrams, SOPs"],
        ["Training", "Administrator and analyst training"],
    ]
    add_table(slide, scope_data, top=1.6, row_height=0.45)

    # Out of scope
    oos_label = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.6), Inches(11.8), Inches(0.4)
    )
    tf = oos_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Out of Scope"
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    oos_points = [
        {"text": "Custom application development"},
        {"text": "Hardware procurement"},
        {"text": "Third-party tool licensing"},
        {"text": "Ongoing managed services (post-warranty)"},
    ]
    add_bullet_points(slide, oos_points, top=5.0, height=2, font_size=14)

    # Deliverables Matrix slide
    slide = create_content_slide(prs, "Deliverables Matrix")
    deliverables_data = [
        ["#", "Deliverable", "Format", "Owner"],
        ["D1", "Solution Architecture Document", "PDF", "Seekurity"],
        ["D2", "Deployment Runbook", "PDF/Wiki", "Seekurity"],
        ["D3", "Log Source Integration Guide", "PDF", "Seekurity"],
        ["D4", "Use Case Documentation", "Excel/PDF", "Seekurity"],
        ["D5", "Dashboard & Report Templates", "Platform", "Seekurity"],
        ["D6", "Training Materials", "PDF/Video", "Seekurity"],
        ["D7", "Knowledge Transfer Sessions", "Live", "Seekurity"],
        ["D8", "Project Completion Report", "PDF", "Seekurity"],
    ]
    add_table(slide, deliverables_data, row_height=0.5)

    # Success Criteria slide
    slide = create_content_slide(prs, "Success Criteria")
    criteria_data = [
        ["Criteria", "Measurement", "Target"],
        [
            "Log source integration",
            "Sources successfully ingesting",
            "100% of agreed sources",
        ],
        ["Data normalization", "Events parsed correctly", ">95% parse success rate"],
        ["Detection coverage", "Use cases deployed", "100% of agreed use cases"],
        ["Alert accuracy", "False positive rate", "<20% after tuning"],
        ["System availability", "Platform uptime", ">99.5% during FVT"],
        ["Knowledge transfer", "Training completion", "100% of designated staff"],
        ["Documentation", "Deliverables accepted", "All D1-D8 signed off"],
    ]
    add_table(slide, criteria_data, row_height=0.5)

    # =========================================================================
    # SECTION 3: TECHNICAL ARCHITECTURE
    # =========================================================================
    create_section_divider(prs, 3, "Technical Architecture")

    # High-Level Architecture slide
    slide = create_content_slide(prs, "High-Level Architecture")
    arch_note = slide.shapes.add_textbox(
        Inches(0.75), Inches(1.6), Inches(11.8), Inches(5)
    )
    tf = arch_note.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "SIEM Platform Components"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = COLORS["dark_gray"]
    run.font.name = "Segoe UI"

    components = [
        "• Collector Tier - Log collection and forwarding",
        "• Indexer Tier - Data indexing and storage",
        "• Search Head - Query and visualization",
        "",
        "Data Sources:",
        "• Firewall & Network Security",
        "• Endpoint Detection & Response (EDR)",
        "• Identity & Access Management (IAM)",
        "• Cloud Infrastructure",
        "• Email & Collaboration",
    ]
    for comp in components:
        add_paragraph(
            tf, comp, font_size=16, font_color=COLORS["dark_gray"], space_before=8
        )

    add_paragraph(tf, "", font_size=12, space_before=12)
    add_paragraph(
        tf,
        "Detailed architecture diagram available in Solution Architecture Document",
        font_size=12,
        font_color=COLORS["light_gray"],
    )

    # Infrastructure Requirements slide
    slide = create_content_slide(
        prs, "Infrastructure Requirements", "Compute Resources"
    )
    infra_data = [
        ["Component", "Specification", "Quantity"],
        ["Collector Nodes", "8 vCPU, 32GB RAM, 500GB SSD", "[X]"],
        ["Indexer Nodes", "16 vCPU, 64GB RAM, 2TB SSD", "[X]"],
        ["Search Heads", "8 vCPU, 32GB RAM, 200GB SSD", "[X]"],
    ]
    add_table(slide, infra_data, top=1.6, row_height=0.5)

    # Network requirements section
    net_label = slide.shapes.add_textbox(
        Inches(0.75), Inches(3.7), Inches(11.8), Inches(0.4)
    )
    tf = net_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Network Requirements"
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    net_data = [
        ["Requirement", "Details"],
        ["Bandwidth", "Minimum [X] Gbps between tiers"],
        ["Latency", "<10ms between components"],
        ["Firewall Rules", "See Firewall Policy Document"],
    ]
    add_table(slide, net_data, top=4.2, row_height=0.45)

    # Log Source Summary slide
    slide = create_content_slide(prs, "Log Source Summary")
    log_data = [
        ["Category", "Sources", "Est. EPS", "Priority"],
        ["Network Security", "Firewall, IDS/IPS, Proxy", "[X]", "P1"],
        ["Endpoint", "EDR, AV, OS Logs", "[X]", "P1"],
        ["Identity", "AD, Azure AD, PAM", "[X]", "P1"],
        ["Cloud", "AWS/Azure/GCP", "[X]", "P2"],
        ["Application", "Web Apps, Databases", "[X]", "P2"],
        ["Email", "O365, Exchange", "[X]", "P2"],
    ]
    add_table(slide, log_data, row_height=0.5)

    # Totals
    totals_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(5.2), Inches(11.8), Inches(0.8)
    )
    tf = totals_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Total Estimated EPS: [X]"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = COLORS["dark_gray"]
    run.font.name = "Segoe UI"
    add_paragraph(
        tf,
        "Total Estimated Daily Volume: [X] GB",
        font_size=14,
        font_bold=True,
        font_color=COLORS["dark_gray"],
        space_before=4,
    )

    # =========================================================================
    # SECTION 4: IMPLEMENTATION APPROACH
    # =========================================================================
    create_section_divider(prs, 4, "Implementation Approach")

    # Implementation Methodology slide
    slide = create_content_slide(
        prs, "Implementation Methodology", "Phase-Based Delivery"
    )
    phases_data = [
        ["Phase", "Name", "Focus"],
        ["Phase 1", "Foundation", "Deploy platform, validate infrastructure"],
        ["Phase 2", "Integration", "Connect log sources, validate parsing"],
        ["Phase 3", "Detection", "Deploy use cases, tune alerts"],
        ["Phase 4", "Transition", "Train team, handover documentation"],
    ]
    add_table(slide, phases_data, row_height=0.55)

    # Phase 1 slide
    slide = create_content_slide(prs, "Phase 1: Foundation", "Duration: [X] weeks")
    p1_data = [
        ["Activity", "Description", "Owner"],
        ["Environment setup", "Deploy infrastructure, base config", "Seekurity"],
        ["Network validation", "Verify connectivity, firewall rules", "Joint"],
        ["Platform installation", "Install and configure SIEM", "Seekurity"],
        ["Health verification", "Validate platform health", "Seekurity"],
        ["Access provisioning", "Create accounts, set permissions", "Joint"],
    ]
    add_table(slide, p1_data, row_height=0.5)

    exit_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(5.0), Inches(11.8), Inches(1.2)
    )
    tf = exit_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Exit Criteria:"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = COLORS["primary_blue"]
    run.font.name = "Segoe UI"
    add_paragraph(
        tf,
        "• Platform deployed and accessible",
        font_size=12,
        font_color=COLORS["dark_gray"],
        space_before=4,
    )
    add_paragraph(
        tf,
        "• All health checks passing",
        font_size=12,
        font_color=COLORS["dark_gray"],
        space_before=2,
    )
    add_paragraph(
        tf,
        "• Administrative access verified",
        font_size=12,
        font_color=COLORS["dark_gray"],
        space_before=2,
    )

    # Phase 2 slide
    slide = create_content_slide(prs, "Phase 2: Integration", "Duration: [X] weeks")
    p2_data = [
        ["Activity", "Description", "Owner"],
        ["P1 source integration", "Connect priority 1 log sources", "Seekurity"],
        ["Parser development", "Create/customize parsers", "Seekurity"],
        ["Data validation", "Verify log ingestion and parsing", "Joint"],
        ["P2 source integration", "Connect priority 2 log sources", "Seekurity"],
        ["Normalization tuning", "Optimize field extraction", "Seekurity"],
    ]
    add_table(slide, p2_data, row_height=0.5)

    exit_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(5.0), Inches(11.8), Inches(1.2)
    )
    tf = exit_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Exit Criteria:"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = COLORS["primary_blue"]
    run.font.name = "Segoe UI"
    add_paragraph(
        tf,
        "• All log sources connected",
        font_size=12,
        font_color=COLORS["dark_gray"],
        space_before=4,
    )
    add_paragraph(
        tf,
        "• >95% parse success rate achieved",
        font_size=12,
        font_color=COLORS["dark_gray"],
        space_before=2,
    )
    add_paragraph(
        tf,
        "• Data flowing to dashboards",
        font_size=12,
        font_color=COLORS["dark_gray"],
        space_before=2,
    )

    # Phase 3 slide
    slide = create_content_slide(prs, "Phase 3: Detection", "Duration: [X] weeks")
    p3_data = [
        ["Activity", "Description", "Owner"],
        ["Use case deployment", "Implement detection rules", "Seekurity"],
        ["Alert configuration", "Set thresholds, routing", "Seekurity"],
        ["Dashboard creation", "Build operational dashboards", "Seekurity"],
        ["Tuning & optimization", "Reduce false positives", "Joint"],
        ["Integration testing", "Validate alert workflow", "Joint"],
    ]
    add_table(slide, p3_data, row_height=0.5)

    exit_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(5.0), Inches(11.8), Inches(1.2)
    )
    tf = exit_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Exit Criteria:"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = COLORS["primary_blue"]
    run.font.name = "Segoe UI"
    add_paragraph(
        tf,
        "• All use cases deployed",
        font_size=12,
        font_color=COLORS["dark_gray"],
        space_before=4,
    )
    add_paragraph(
        tf,
        "• False positive rate <20%",
        font_size=12,
        font_color=COLORS["dark_gray"],
        space_before=2,
    )
    add_paragraph(
        tf,
        "• Alerts routing correctly",
        font_size=12,
        font_color=COLORS["dark_gray"],
        space_before=2,
    )

    # Phase 4 slide
    slide = create_content_slide(prs, "Phase 4: Transition", "Duration: [X] weeks")
    p4_data = [
        ["Activity", "Description", "Owner"],
        ["Documentation delivery", "Finalize all deliverables", "Seekurity"],
        ["Administrator training", "Platform management training", "Seekurity"],
        ["Analyst training", "Detection and response training", "Seekurity"],
        ["FVT execution", "User acceptance testing", customer_name],
        ["Project closure", "Sign-off, transition to support", "Joint"],
    ]
    add_table(slide, p4_data, row_height=0.5)

    exit_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(5.0), Inches(11.8), Inches(1.2)
    )
    tf = exit_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Exit Criteria:"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = COLORS["primary_blue"]
    run.font.name = "Segoe UI"
    add_paragraph(
        tf,
        "• All deliverables accepted",
        font_size=12,
        font_color=COLORS["dark_gray"],
        space_before=4,
    )
    add_paragraph(
        tf,
        "• Training completed",
        font_size=12,
        font_color=COLORS["dark_gray"],
        space_before=2,
    )
    add_paragraph(
        tf,
        "• FVT passed and sign-off received",
        font_size=12,
        font_color=COLORS["dark_gray"],
        space_before=2,
    )

    # =========================================================================
    # SECTION 5: TIMELINE & MILESTONES
    # =========================================================================
    create_section_divider(prs, 5, "Timeline & Milestones")

    # Project Timeline slide
    slide = create_content_slide(prs, "Project Timeline")
    timeline_data = [
        ["Phase", "Start", "End", "Duration"],
        ["Phase 1: Foundation", "[DATE]", "[DATE]", "[X] weeks"],
        ["Phase 2: Integration", "[DATE]", "[DATE]", "[X] weeks"],
        ["Phase 3: Detection", "[DATE]", "[DATE]", "[X] weeks"],
        ["Phase 4: Transition", "[DATE]", "[DATE]", "[X] weeks"],
        ["Project Complete", "-", "[DATE]", "[X] weeks total"],
    ]
    add_table(slide, timeline_data, row_height=0.55)

    # Key Milestones slide
    slide = create_content_slide(prs, "Key Milestones")
    milestones_data = [
        ["#", "Milestone", "Target Date", "Dependencies"],
        ["M1", "Infrastructure Ready", "[DATE]", "Environment provisioned"],
        ["M2", "Platform Deployed", "[DATE]", "M1 complete"],
        ["M3", "P1 Sources Integrated", "[DATE]", "Network access, credentials"],
        ["M4", "All Sources Integrated", "[DATE]", "M3 complete"],
        ["M5", "Use Cases Deployed", "[DATE]", "M4 complete"],
        ["M6", "Training Complete", "[DATE]", "M5 complete"],
        ["M7", "FVT Sign-off", "[DATE]", "M6 complete"],
        ["M8", "Project Closure", "[DATE]", "M7 complete"],
    ]
    add_table(slide, milestones_data, row_height=0.48)

    # Dependencies slide
    slide = create_content_slide(prs, "Dependencies & Assumptions", "Key Dependencies")
    deps_data = [
        ["#", "Dependency", "Owner", "Required By"],
        ["D1", "Infrastructure provisioned", customer_name, "[DATE]"],
        ["D2", "Network connectivity established", customer_name, "[DATE]"],
        ["D3", "Firewall rules implemented", customer_name, "[DATE]"],
        ["D4", "Log source credentials provided", customer_name, "[DATE]"],
        ["D5", "Technical resources available", customer_name, "Ongoing"],
    ]
    add_table(slide, deps_data, top=1.6, row_height=0.5)

    # Assumptions
    assume_label = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.4), Inches(11.8), Inches(0.4)
    )
    tf = assume_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Key Assumptions"
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    assumptions = [
        {"text": "Customer resources available per agreed schedule"},
        {"text": "No major infrastructure changes during implementation"},
        {"text": "Log sources accessible from SIEM platform"},
        {"text": "Change management approvals obtained timely"},
    ]
    add_bullet_points(slide, assumptions, top=4.9, height=2, font_size=14)

    # =========================================================================
    # SECTION 6: GOVERNANCE & COMMUNICATION
    # =========================================================================
    create_section_divider(prs, 6, "Governance & Communication")

    # Governance Model slide
    slide = create_content_slide(prs, "Governance Model", "Decision Authority")
    gov_data = [
        ["Decision Type", "Authority", "Escalation Path"],
        ["Day-to-day technical", "Technical Leads", "Project Managers"],
        ["Scope changes", "Project Managers", "Project Sponsors"],
        ["Schedule changes", "Project Managers", "Project Sponsors"],
        ["Resource allocation", "Project Managers", "Project Sponsors"],
        ["Contract changes", "Project Sponsors", "Executive Leadership"],
    ]
    add_table(slide, gov_data, top=1.6, row_height=0.5)

    # Change Control note
    cc_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.8), Inches(11.8), Inches(1.5)
    )
    tf = cc_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Change Control"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = COLORS["primary_blue"]
    run.font.name = "Segoe UI"
    add_paragraph(
        tf,
        "All scope, timeline, or budget changes require:",
        font_size=12,
        font_color=COLORS["dark_gray"],
        space_before=8,
    )
    add_paragraph(
        tf,
        "1. Written change request submission    2. Impact assessment",
        font_size=12,
        font_color=COLORS["dark_gray"],
        space_before=4,
    )
    add_paragraph(
        tf,
        "3. Approval from both Project Managers    4. Documentation update",
        font_size=12,
        font_color=COLORS["dark_gray"],
        space_before=4,
    )

    # Communication Plan slide
    slide = create_content_slide(prs, "Communication Plan")
    comm_data = [
        ["Meeting", "Frequency", "Attendees", "Purpose"],
        ["Daily Standup", "Daily", "Technical team", "Progress, blockers"],
        ["Weekly Status", "Weekly", "PM + Leads", "Status review, risks"],
        ["Steering Committee", "Bi-weekly", "Sponsors + PMs", "Decisions, escalations"],
        ["Technical Review", "As needed", "Technical team", "Deep-dive sessions"],
    ]
    add_table(slide, comm_data, top=1.6, row_height=0.5)

    # Reporting section
    report_label = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.2), Inches(11.8), Inches(0.4)
    )
    tf = report_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Reporting"
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    report_data = [
        ["Report", "Frequency", "Distribution"],
        ["Status Report", "Weekly", "All stakeholders"],
        ["Risk Register", "Weekly", "PMs, Sponsors"],
        ["Milestone Report", "Per milestone", "All stakeholders"],
    ]
    add_table(slide, report_data, top=4.7, row_height=0.45, font_size=11)

    # Key Contacts slide
    slide = create_content_slide(prs, "Communication Channels")
    channels_data = [
        ["Channel", "Purpose", "Response Time"],
        ["Email", "Formal communication, documentation", "24 hours"],
        ["[Collaboration Tool]", "Day-to-day coordination", "4 hours"],
        ["Phone/Video", "Urgent issues, meetings", "Immediate"],
        ["Ticketing System", "Issue tracking, requests", "Per SLA"],
    ]
    add_table(slide, channels_data, top=1.6, row_height=0.5)

    # Key contacts section
    contacts_label = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.0), Inches(11.8), Inches(0.4)
    )
    tf = contacts_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Key Contacts"
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    contacts_data = [
        ["Role", "Name", "Email", "Phone"],
        ["Seekurity PM", "[NAME]", "[EMAIL]", "[PHONE]"],
        ["Seekurity Tech Lead", "[NAME]", "[EMAIL]", "[PHONE]"],
        [f"{customer_name} PM", "[NAME]", "[EMAIL]", "[PHONE]"],
        [f"{customer_name} Tech Lead", "[NAME]", "[EMAIL]", "[PHONE]"],
    ]
    add_table(slide, contacts_data, top=4.5, row_height=0.45, font_size=11)

    # =========================================================================
    # SECTION 7: RISKS & MITIGATION
    # =========================================================================
    create_section_divider(prs, 7, "Risks & Mitigation")

    # Risk Register slide
    slide = create_content_slide(prs, "Risk Register")
    risk_data = [
        ["#", "Risk", "Probability", "Impact", "Mitigation"],
        [
            "R1",
            "Infrastructure delays",
            "Medium",
            "High",
            "Early engagement, weekly tracking",
        ],
        [
            "R2",
            "Resource unavailability",
            "Medium",
            "Medium",
            "Backup resources identified",
        ],
        [
            "R3",
            "Log source access issues",
            "High",
            "Medium",
            "Pre-validation checklist",
        ],
        ["R4", "Scope creep", "Medium", "High", "Strict change control"],
        ["R5", "Integration complexity", "Medium", "Medium", "Technical POC early"],
        [
            "R6",
            "Data volume exceeds estimate",
            "Low",
            "High",
            "Capacity planning buffer",
        ],
    ]
    add_table(slide, risk_data, row_height=0.55, font_size=11)

    # Risk Response Actions slide
    slide = create_content_slide(
        prs, "Risk Response Actions", "Immediate Actions Required"
    )
    actions_data = [
        ["#", "Action", "Owner", "Due Date"],
        ["A1", "Confirm infrastructure timeline", f"{customer_name} PM", "[DATE]"],
        ["A2", "Validate network connectivity", "Technical Leads", "[DATE]"],
        ["A3", "Obtain log source credentials", f"{customer_name} Tech", "[DATE]"],
        ["A4", "Schedule resource availability", "Both PMs", "[DATE]"],
    ]
    add_table(slide, actions_data, top=1.6, row_height=0.5)

    # Escalation triggers
    esc_label = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.0), Inches(11.8), Inches(0.4)
    )
    tf = esc_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Escalation Triggers"
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    triggers = [
        {"text": 'Any risk probability increasing to "High"'},
        {"text": "Any milestone at risk of delay >5 days"},
        {"text": "Resource conflicts lasting >3 days"},
        {"text": "Unresolved blockers >48 hours"},
    ]
    add_bullet_points(slide, triggers, top=4.5, height=2, font_size=14)

    # =========================================================================
    # SECTION 8: NEXT STEPS & Q&A
    # =========================================================================
    create_section_divider(prs, 8, "Next Steps & Q&A")

    # Immediate Next Steps slide
    slide = create_content_slide(prs, "Immediate Next Steps")
    next_steps_data = [
        ["#", "Action", "Owner", "Due Date"],
        ["1", "Distribute meeting minutes", "Seekurity PM", "[DATE]"],
        ["2", "Finalize infrastructure requirements", customer_name, "[DATE]"],
        ["3", "Submit firewall change requests", customer_name, "[DATE]"],
        ["4", "Provide log source credentials", customer_name, "[DATE]"],
        ["5", "Schedule Phase 1 kickoff", "Both PMs", "[DATE]"],
        ["6", "Set up collaboration channels", "Both PMs", "[DATE]"],
    ]
    add_table(slide, next_steps_data, row_height=0.55)

    # Q&A slide
    slide = create_content_slide(prs, "Questions & Discussion", "Topics for Discussion")
    qa_points = [
        {"text": "Infrastructure readiness timeline"},
        {"text": "Resource availability confirmation"},
        {"text": "Log source access verification"},
        {"text": "Communication tool preferences"},
        {"text": "Any concerns or clarifications"},
    ]
    add_bullet_points(slide, qa_points, font_size=20)

    # =========================================================================
    # THANK YOU SLIDE
    # =========================================================================
    create_thank_you_slide(prs)

    # =========================================================================
    # APPENDIX
    # =========================================================================
    slide = create_content_slide(prs, "Appendix", "Document References")
    appendix_data = [
        ["Document", "Location"],
        ["Statement of Work", "[LINK]"],
        ["Solution Architecture", "[LINK]"],
        ["Firewall Policy", "[LINK]"],
        ["Log Source Matrix", "[LINK]"],
        ["Deployment Schedule", "[LINK]"],
    ]
    add_table(slide, appendix_data, top=1.6, row_height=0.5)

    # Version history
    ver_label = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.8), Inches(11.8), Inches(0.4)
    )
    tf = ver_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Version History"
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    ver_data = [
        ["Version", "Date", "Author", "Changes"],
        ["1.0", "[DATE]", "[NAME]", "Initial version"],
    ]
    add_table(slide, ver_data, top=5.3, row_height=0.45, font_size=11)

    # Save presentation
    prs.save(output_path)
    print(f"Kickoff presentation saved to: {output_path}")
    return prs


# =============================================================================
# COMPLETION PRESENTATION GENERATOR
# =============================================================================


def generate_completion_presentation(
    output_path, customer_name="[CUSTOMER_NAME]", date="[DATE]"
):
    """Generate the Project Completion presentation."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # =========================================================================
    # TITLE SLIDE
    # =========================================================================
    create_title_slide(
        prs, "SIEM Implementation", "Project Completion Report", customer_name, date
    )

    # =========================================================================
    # AGENDA
    # =========================================================================
    slide = create_content_slide(prs, "Agenda")
    agenda_data = [
        ["#", "Topic", "Duration"],
        ["1", "Executive Summary", "10 min"],
        ["2", "Project Achievements", "15 min"],
        ["3", "Technical Delivery", "20 min"],
        ["4", "Metrics & Performance", "10 min"],
        ["5", "Lessons Learned", "10 min"],
        ["6", "Support & Maintenance", "10 min"],
        ["7", "Recommendations", "10 min"],
        ["8", "Sign-off & Closure", "5 min"],
    ]
    add_table(slide, agenda_data)

    note_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(5.8), Inches(11.8), Inches(0.4)
    )
    tf = note_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Total Duration: ~90 minutes"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    # =========================================================================
    # SECTION 1: EXECUTIVE SUMMARY
    # =========================================================================
    create_section_divider(prs, 1, "Executive Summary")

    # Project Overview slide
    slide = create_content_slide(prs, "Project Overview", "Engagement Summary")
    overview_data = [
        ["Attribute", "Detail"],
        ["Customer", customer_name],
        ["Project", "SIEM Implementation"],
        ["Start Date", "[START_DATE]"],
        ["End Date", "[END_DATE]"],
        ["Duration", "[X] weeks"],
        ["Project Manager", "[PM_NAME]"],
        ["Technical Lead", "[TECH_LEAD_NAME]"],
    ]
    add_table(slide, overview_data, row_height=0.5)

    # Key Outcomes slide
    slide = create_content_slide(prs, "Key Outcomes", "Project Status: COMPLETED")
    outcomes_data = [
        ["Metric", "Target", "Achieved", "Status"],
        ["Log Sources Integrated", "[X]", "[X]", "100%"],
        ["Use Cases Deployed", "[X]", "[X]", "100%"],
        ["Training Sessions", "[X]", "[X]", "100%"],
        ["Documentation Deliverables", "[X]", "[X]", "100%"],
        ["On-Time Delivery", "[DATE]", "[DATE]", "On Schedule"],
        ["Budget", "$[X]", "$[X]", "On Budget"],
    ]
    add_table(slide, outcomes_data, row_height=0.5)

    # Value Delivered slide
    slide = create_content_slide(
        prs, "Value Delivered", "Security Capabilities Enabled"
    )
    value_data = [
        ["Capability", "Before", "After"],
        ["Log Visibility", "[X] sources", "[X] sources"],
        ["Detection Coverage", "[X]%", "[X]%"],
        ["Mean Time to Detect", "[X] hours", "[X] minutes"],
        ["Compliance Reporting", "Manual", "Automated"],
        ["Incident Response", "Ad-hoc", "Structured"],
    ]
    add_table(slide, value_data, top=1.6, row_height=0.5)

    # Business Impact
    impact_label = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.5), Inches(11.8), Inches(0.4)
    )
    tf = impact_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Business Impact"
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    impacts = [
        {
            "text": "Reduced Risk: Comprehensive threat visibility across infrastructure",
            "bold": True,
        },
        {
            "text": "Improved Efficiency: Automated detection and correlation",
            "bold": True,
        },
        {"text": "Compliance Ready: Audit-ready logging and reporting", "bold": True},
        {
            "text": "Operational Maturity: Documented processes and trained team",
            "bold": True,
        },
    ]
    add_bullet_points(slide, impacts, top=5.0, height=2, font_size=14)

    # =========================================================================
    # SECTION 2: PROJECT ACHIEVEMENTS
    # =========================================================================
    create_section_divider(prs, 2, "Project Achievements")

    # Milestone Completion slide
    slide = create_content_slide(prs, "Milestone Completion")
    milestones_data = [
        ["#", "Milestone", "Target", "Actual", "Variance"],
        ["M1", "Infrastructure Ready", "[DATE]", "[DATE]", "0 days"],
        ["M2", "Platform Deployed", "[DATE]", "[DATE]", "0 days"],
        ["M3", "P1 Sources Integrated", "[DATE]", "[DATE]", "0 days"],
        ["M4", "All Sources Integrated", "[DATE]", "[DATE]", "[X] days"],
        ["M5", "Use Cases Deployed", "[DATE]", "[DATE]", "0 days"],
        ["M6", "Training Complete", "[DATE]", "[DATE]", "0 days"],
        ["M7", "FVT Sign-off", "[DATE]", "[DATE]", "0 days"],
        ["M8", "Project Closure", "[DATE]", "[DATE]", "0 days"],
    ]
    add_table(slide, milestones_data, row_height=0.48, font_size=11)

    perf_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(6.0), Inches(11.8), Inches(0.4)
    )
    tf = perf_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Overall Schedule Performance: On Time / [X] days early/late"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = COLORS["success_green"]
    run.font.name = "Segoe UI"

    # Deliverables Summary slide
    slide = create_content_slide(prs, "Deliverables Summary")
    deliverables_data = [
        ["#", "Deliverable", "Status", "Accepted"],
        ["D1", "Solution Architecture Document", "Delivered", "[DATE]"],
        ["D2", "Deployment Runbook", "Delivered", "[DATE]"],
        ["D3", "Log Source Integration Guide", "Delivered", "[DATE]"],
        ["D4", "Use Case Documentation", "Delivered", "[DATE]"],
        ["D5", "Dashboard & Report Templates", "Delivered", "[DATE]"],
        ["D6", "Training Materials", "Delivered", "[DATE]"],
        ["D7", "Knowledge Transfer Sessions", "Completed", "[DATE]"],
        ["D8", "Project Completion Report", "This document", "[DATE]"],
    ]
    add_table(slide, deliverables_data, row_height=0.48, font_size=11)

    accept_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(6.0), Inches(11.8), Inches(0.4)
    )
    tf = accept_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"All deliverables accepted by {customer_name}"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = COLORS["success_green"]
    run.font.name = "Segoe UI"

    # Scope Summary slide
    slide = create_content_slide(prs, "Scope Summary", "Delivered as Planned")
    scope_points = [
        {"text": "Platform deployment and configuration"},
        {"text": "[X] log sources integrated and validated"},
        {"text": "[X] detection use cases implemented"},
        {"text": "[X] dashboards and reports created"},
        {"text": "Administrator and analyst training"},
        {"text": "Complete documentation package"},
    ]
    add_bullet_points(slide, scope_points, font_size=16)

    # Change Requests
    cr_label = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.5), Inches(11.8), Inches(0.4)
    )
    tf = cr_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Change Requests Processed"
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    cr_data = [
        ["CR#", "Description", "Impact", "Status"],
        ["CR-001", "[Description]", "[Days/Cost]", "Approved/Completed"],
        ["CR-002", "[Description]", "[Days/Cost]", "Approved/Completed"],
    ]
    add_table(slide, cr_data, top=5.0, row_height=0.45, font_size=11)

    # =========================================================================
    # SECTION 3: TECHNICAL DELIVERY
    # =========================================================================
    create_section_divider(prs, 3, "Technical Delivery")

    # Architecture Deployed slide
    slide = create_content_slide(prs, "Architecture Deployed", "Production Environment")
    arch_data = [
        ["Component", "Specification", "Count", "Status"],
        ["Collector Nodes", "8 vCPU, 32GB RAM", "[X]", "Operational"],
        ["Indexer Nodes", "16 vCPU, 64GB RAM", "[X]", "Operational"],
        ["Search Heads", "8 vCPU, 32GB RAM", "[X]", "Operational"],
        ["Storage", "[X] TB", "[X]", "Provisioned"],
    ]
    add_table(slide, arch_data, top=1.6, row_height=0.5)

    # HA section
    ha_label = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.2), Inches(11.8), Inches(0.4)
    )
    tf = ha_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "High Availability"
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    ha_points = [
        {"text": "Clustering: Enabled"},
        {"text": "Replication Factor: [X]"},
        {"text": "Search Factor: [X]"},
        {"text": "Backup: [Schedule]"},
    ]
    add_bullet_points(slide, ha_points, top=4.7, height=2, font_size=14)

    # Log Source Integration slide
    slide = create_content_slide(prs, "Log Source Integration", "By Category")
    log_data = [
        ["Category", "Planned", "Integrated", "EPS"],
        ["Network Security", "[X]", "[X]", "[X]"],
        ["Endpoint Security", "[X]", "[X]", "[X]"],
        ["Identity & Access", "[X]", "[X]", "[X]"],
        ["Cloud Infrastructure", "[X]", "[X]", "[X]"],
        ["Applications", "[X]", "[X]", "[X]"],
        ["Email & Collaboration", "[X]", "[X]", "[X]"],
        ["Total", "[X]", "[X]", "[X]"],
    ]
    add_table(slide, log_data, row_height=0.52)

    # Detection Use Cases slide
    slide = create_content_slide(
        prs, "Detection Use Cases", "Deployed Use Cases by Category"
    )
    usecase_data = [
        ["Category", "Count", "Examples"],
        ["Authentication", "[X]", "Brute force, impossible travel, MFA bypass"],
        ["Malware", "[X]", "Known IOCs, suspicious execution, C2"],
        ["Data Exfiltration", "[X]", "Large transfers, USB usage, cloud upload"],
        ["Privilege Escalation", "[X]", "Admin creation, permission changes"],
        ["Network Anomaly", "[X]", "Port scanning, lateral movement"],
        ["Compliance", "[X]", "Audit log tampering, policy violations"],
        ["Total", "[X]", ""],
    ]
    add_table(slide, usecase_data, row_height=0.52)

    # Use Case Performance slide
    slide = create_content_slide(
        prs, "Use Case Performance", "Alert Statistics (Last 30 Days)"
    )
    stats_data = [
        ["Metric", "Value"],
        ["Total Alerts Generated", "[X]"],
        ["True Positives", "[X] ([X]%)"],
        ["False Positives", "[X] ([X]%)"],
        ["Average Daily Alerts", "[X]"],
        ["Mean Time to Alert", "[X] seconds"],
    ]
    add_table(slide, stats_data, top=1.6, row_height=0.5)

    # Top triggered
    top_label = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.4), Inches(11.8), Inches(0.4)
    )
    tf = top_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Top Triggered Use Cases"
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    top_data = [
        ["Rank", "Use Case", "Count", "Action Rate"],
        ["1", "[Use Case Name]", "[X]", "[X]%"],
        ["2", "[Use Case Name]", "[X]", "[X]%"],
        ["3", "[Use Case Name]", "[X]", "[X]%"],
    ]
    add_table(slide, top_data, top=4.9, row_height=0.45, font_size=11)

    # Dashboards & Reports slide
    slide = create_content_slide(prs, "Dashboards & Reports", "Operational Dashboards")
    dash_data = [
        ["Dashboard", "Purpose", "Refresh"],
        ["Security Operations Center", "Real-time monitoring", "1 min"],
        ["Executive Summary", "KPIs and trends", "Daily"],
        ["Threat Landscape", "Attack patterns", "Hourly"],
        ["Compliance Status", "Audit readiness", "Daily"],
        ["System Health", "Platform performance", "5 min"],
    ]
    add_table(slide, dash_data, top=1.6, row_height=0.5)

    # Scheduled Reports
    report_label = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.5), Inches(11.8), Inches(0.4)
    )
    tf = report_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Scheduled Reports"
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    report_data = [
        ["Report", "Frequency", "Distribution"],
        ["Daily Security Summary", "Daily 8:00 AM", "SOC Team"],
        ["Weekly Executive Report", "Monday 9:00 AM", "Leadership"],
        ["Monthly Compliance Report", "1st of month", "Compliance Team"],
    ]
    add_table(slide, report_data, top=5.0, row_height=0.45, font_size=11)

    # Integrations slide
    slide = create_content_slide(prs, "Integrations", "Connected Systems")
    int_data = [
        ["System", "Type", "Integration", "Status"],
        ["[Ticketing System]", "ITSM", "REST API", "Active"],
        ["[SOAR Platform]", "Automation", "Webhook", "Active"],
        ["[Threat Intel]", "TI Feed", "API", "Active"],
        ["[Email]", "Notification", "SMTP", "Active"],
    ]
    add_table(slide, int_data, top=1.6, row_height=0.5)

    # Workflow Automation
    wf_label = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.0), Inches(11.8), Inches(0.4)
    )
    tf = wf_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Workflow Automation"
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    wf_data = [
        ["Workflow", "Trigger", "Action"],
        ["Critical Alert", "Severity = Critical", "Create P1 ticket, notify on-call"],
        ["Malware Detection", "Malware use case", "Isolate endpoint (SOAR)"],
        ["User Compromise", "Auth anomaly", "Disable account, notify"],
    ]
    add_table(slide, wf_data, top=4.5, row_height=0.45, font_size=11)

    # =========================================================================
    # SECTION 4: METRICS & PERFORMANCE
    # =========================================================================
    create_section_divider(prs, 4, "Metrics & Performance")

    # Platform Performance slide
    slide = create_content_slide(prs, "Platform Performance", "System Health")
    perf_data = [
        ["Metric", "Target", "Current", "Status"],
        ["Platform Uptime", ">99.5%", "[X]%", "Pass"],
        ["Search Response Time", "<5 sec", "[X] sec", "Pass"],
        ["Indexing Latency", "<30 sec", "[X] sec", "Pass"],
        ["Storage Utilization", "<80%", "[X]%", "Pass"],
        ["CPU Utilization (avg)", "<70%", "[X]%", "Pass"],
        ["Memory Utilization", "<85%", "[X]%", "Pass"],
    ]
    add_table(slide, perf_data, row_height=0.55)

    # Data Metrics slide
    slide = create_content_slide(prs, "Data Metrics", "Ingestion Statistics")
    ingest_data = [
        ["Metric", "Value"],
        ["Average Daily Volume", "[X] GB"],
        ["Peak Daily Volume", "[X] GB"],
        ["Average EPS", "[X]"],
        ["Peak EPS", "[X]"],
        ["Total Events (Project Duration)", "[X] billion"],
        ["Parse Success Rate", "[X]%"],
    ]
    add_table(slide, ingest_data, top=1.6, row_height=0.5)

    # Storage projection
    storage_label = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.8), Inches(11.8), Inches(0.4)
    )
    tf = storage_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Storage Projection"
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    storage_data = [
        ["Retention", "Current Usage", "Projected (1 year)"],
        ["Hot", "[X] TB", "[X] TB"],
        ["Warm", "[X] TB", "[X] TB"],
        ["Cold", "[X] TB", "[X] TB"],
    ]
    add_table(slide, storage_data, top=5.3, row_height=0.4, font_size=11)

    # Security Metrics slide
    slide = create_content_slide(prs, "Security Metrics", "Detection Effectiveness")
    sec_data = [
        ["Metric", "Baseline", "Current", "Improvement"],
        ["Detection Coverage", "[X]%", "[X]%", "+[X]%"],
        ["MTTD (Mean Time to Detect)", "[X] hrs", "[X] min", "[X]% faster"],
        ["MTTR (Mean Time to Respond)", "[X] hrs", "[X] hrs", "[X]% faster"],
        ["Alert Fidelity", "N/A", "[X]%", "New capability"],
    ]
    add_table(slide, sec_data, top=1.6, row_height=0.55)

    # Compliance
    comp_label = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.2), Inches(11.8), Inches(0.4)
    )
    tf = comp_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Compliance Readiness"
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    comp_data = [
        ["Framework", "Logging Coverage", "Report Automation"],
        ["[Framework 1]", "[X]%", "Automated"],
        ["[Framework 2]", "[X]%", "Automated"],
        ["[Framework 3]", "[X]%", "Automated"],
    ]
    add_table(slide, comp_data, top=4.7, row_height=0.45, font_size=11)

    # =========================================================================
    # SECTION 5: LESSONS LEARNED
    # =========================================================================
    create_section_divider(prs, 5, "Lessons Learned")

    # What Went Well slide
    slide = create_content_slide(prs, "What Went Well")
    well_data = [
        ["Category", "Observation"],
        ["Planning", "Thorough scoping minimized scope creep"],
        ["Communication", "Regular status meetings kept all parties aligned"],
        ["Technical", "Pre-validation checklist reduced integration issues"],
        ["Collaboration", "Strong customer engagement accelerated decisions"],
        ["Quality", "Iterative testing caught issues early"],
        ["Documentation", "Comprehensive docs enabled smooth handover"],
    ]
    add_table(slide, well_data, row_height=0.55)

    # Areas for Improvement slide
    slide = create_content_slide(prs, "Areas for Improvement")
    improve_data = [
        ["Category", "Observation", "Recommendation"],
        [
            "Access",
            "Credential provisioning took longer",
            "Earlier engagement with IAM team",
        ],
        ["Testing", "FVT window was tight", "Allocate buffer for FVT"],
        ["Training", "Some attendees missed sessions", "Record sessions for on-demand"],
        ["Change Control", "Some requests lacked detail", "Improve CR template"],
    ]
    add_table(slide, improve_data, row_height=0.6)

    # Recommendations slide
    slide = create_content_slide(prs, "Recommendations for Future Projects")

    # Process improvements
    process_label = slide.shapes.add_textbox(
        Inches(0.75), Inches(1.6), Inches(11.8), Inches(0.4)
    )
    tf = process_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Process Improvements"
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = COLORS["primary_blue"]
    run.font.name = "Segoe UI"

    process_points = [
        {
            "text": "Pre-engagement checklist: Validate access and credentials before kickoff"
        },
        {"text": "Stakeholder mapping: Identify all required approvers early"},
        {"text": "Training scheduling: Book training dates at project start"},
        {"text": "Buffer allocation: Include contingency in critical path activities"},
    ]
    add_bullet_points(slide, process_points, top=2.1, height=2, font_size=14)

    # Technical improvements
    tech_label = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.2), Inches(11.8), Inches(0.4)
    )
    tf = tech_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Technical Improvements"
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = COLORS["primary_blue"]
    run.font.name = "Segoe UI"

    tech_points = [
        {"text": "Parser library: Maintain reusable parser templates"},
        {"text": "Use case catalog: Standardize detection content"},
        {"text": "Integration patterns: Document common integration approaches"},
    ]
    add_bullet_points(slide, tech_points, top=4.7, height=2, font_size=14)

    # =========================================================================
    # SECTION 6: SUPPORT & MAINTENANCE
    # =========================================================================
    create_section_divider(prs, 6, "Support & Maintenance")

    # Support Transition slide
    slide = create_content_slide(prs, "Support Transition", "Warranty Period")
    warranty_data = [
        ["Attribute", "Detail"],
        ["Start Date", "[DATE]"],
        ["End Date", "[DATE]"],
        ["Duration", "[X] days"],
        ["Coverage", "[Coverage details]"],
    ]
    add_table(slide, warranty_data, top=1.6, row_height=0.5)

    # Support channels
    support_label = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.2), Inches(11.8), Inches(0.4)
    )
    tf = support_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Support Channels"
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    support_data = [
        ["Channel", "Contact", "Hours"],
        ["Email", "[support@seekurity.com]", "24x7"],
        ["Phone", "[PHONE]", "Business hours"],
        ["Portal", "[URL]", "24x7"],
    ]
    add_table(slide, support_data, top=4.7, row_height=0.45, font_size=11)

    # Support SLAs slide
    slide = create_content_slide(prs, "Support SLAs", "Response Times")
    sla_data = [
        ["Severity", "Description", "Response", "Resolution"],
        ["P1 - Critical", "Platform down, no workaround", "1 hour", "4 hours"],
        ["P2 - High", "Major feature impacted", "4 hours", "8 hours"],
        ["P3 - Medium", "Minor feature impacted", "8 hours", "24 hours"],
        ["P4 - Low", "General questions", "24 hours", "72 hours"],
    ]
    add_table(slide, sla_data, top=1.6, row_height=0.55)

    # Escalation path
    esc_label = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.4), Inches(11.8), Inches(0.4)
    )
    tf = esc_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Escalation Path"
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    esc_data = [
        ["Level", "Contact", "Timeframe"],
        ["L1", "Support Engineer", "Initial"],
        ["L2", "Senior Engineer", "P1: 2hr, P2: 4hr"],
        ["L3", "Technical Lead", "P1: 4hr, P2: 8hr"],
        ["Management", "Support Manager", "P1: 8hr"],
    ]
    add_table(slide, esc_data, top=4.9, row_height=0.4, font_size=11)

    # Ongoing Maintenance slide
    slide = create_content_slide(
        prs, "Ongoing Maintenance", "Customer Responsibilities"
    )
    maint_data = [
        ["Activity", "Frequency", "Responsibility"],
        ["Platform health monitoring", "Daily", "SOC Team"],
        ["Use case tuning", "Weekly", "Security Analyst"],
        ["Storage management", "Monthly", "Platform Admin"],
        ["Backup verification", "Monthly", "Platform Admin"],
        ["Version updates", "Quarterly", "Platform Admin"],
        ["Use case review", "Quarterly", "Security Team"],
    ]
    add_table(slide, maint_data, row_height=0.55)

    # Knowledge Transfer slide
    slide = create_content_slide(
        prs, "Knowledge Transfer Summary", "Training Completed"
    )
    kt_data = [
        ["Session", "Attendees", "Duration", "Date"],
        ["Platform Administration", "[X]", "[X] hours", "[DATE]"],
        ["Security Analysis", "[X]", "[X] hours", "[DATE]"],
        ["Use Case Management", "[X]", "[X] hours", "[DATE]"],
        ["Report Development", "[X]", "[X] hours", "[DATE]"],
    ]
    add_table(slide, kt_data, top=1.6, row_height=0.5)

    # Materials provided
    mat_label = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.2), Inches(11.8), Inches(0.4)
    )
    tf = mat_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Materials Provided"
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    materials = [
        {"text": "Administrator Guide"},
        {"text": "Analyst Playbook"},
        {"text": "Quick Reference Cards"},
        {"text": "Video Recordings"},
        {"text": "Lab Exercises"},
    ]
    add_bullet_points(slide, materials, top=4.7, height=2, font_size=14)

    # =========================================================================
    # SECTION 7: RECOMMENDATIONS
    # =========================================================================
    create_section_divider(prs, 7, "Recommendations")

    # Short-Term Recommendations slide
    slide = create_content_slide(prs, "Short-Term Recommendations", "Next 90 Days")
    short_data = [
        ["#", "Recommendation", "Priority", "Effort"],
        ["1", "Complete tuning of all use cases", "High", "Medium"],
        ["2", "Develop SOC playbooks for top alerts", "High", "Medium"],
        ["3", "Integrate additional threat intel feeds", "Medium", "Low"],
        ["4", "Create custom compliance reports", "Medium", "Medium"],
        ["5", "Document incident response procedures", "High", "Medium"],
    ]
    add_table(slide, short_data, row_height=0.55)

    # Medium-Term Recommendations slide
    slide = create_content_slide(prs, "Medium-Term Recommendations", "3-12 Months")
    medium_data = [
        ["#", "Recommendation", "Priority", "Effort"],
        ["1", "Implement SOAR automation for top use cases", "High", "High"],
        ["2", "Expand log source coverage (Phase 2)", "Medium", "High"],
        ["3", "Develop advanced detection content", "Medium", "Medium"],
        ["4", "Implement user behavior analytics", "Medium", "High"],
        ["5", "Conduct purple team exercises", "Medium", "Medium"],
        ["6", "Achieve 24x7 SOC coverage", "High", "High"],
    ]
    add_table(slide, medium_data, row_height=0.52)

    # Long-Term Roadmap slide
    slide = create_content_slide(prs, "Long-Term Roadmap", "12+ Months")
    long_data = [
        ["Initiative", "Description", "Business Value"],
        ["Threat Hunting Program", "Proactive threat discovery", "Reduced dwell time"],
        ["ML-Based Detection", "Anomaly detection at scale", "Novel threat detection"],
        ["Security Data Lake", "Long-term analytics", "Forensic capability"],
        ["Multi-Cloud Coverage", "Unified cloud visibility", "Cloud security posture"],
        ["XDR Integration", "Extended detection", "Correlated response"],
    ]
    add_table(slide, long_data, row_height=0.6)

    # Investment Recommendations slide
    slide = create_content_slide(
        prs, "Investment Recommendations", "Estimated Resource Requirements"
    )
    invest_data = [
        ["Initiative", "Type", "Estimated Investment"],
        ["SOAR Implementation", "Technology + Services", "$[X]"],
        ["Phase 2 Log Sources", "Services", "$[X]"],
        ["Advanced Detection Content", "Services", "$[X]"],
        ["Managed Detection Service", "Annual Service", "$[X]/year"],
    ]
    add_table(slide, invest_data, top=1.6, row_height=0.55)

    # Services available
    svc_label = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.4), Inches(11.8), Inches(0.4)
    )
    tf = svc_label.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Seekurity Services Available"
    run.font.size = Pt(20)
    run.font.bold = False
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    services = [
        {"text": "Managed Detection & Response (MDR)"},
        {"text": "Security Operations Center as a Service (SOCaaS)"},
        {"text": "Threat Hunting Services"},
        {"text": "Purple Team Exercises"},
        {"text": "SIEM Optimization Services"},
    ]
    add_bullet_points(slide, services, top=4.9, height=2, font_size=14)

    # =========================================================================
    # SECTION 8: SIGN-OFF & CLOSURE
    # =========================================================================
    create_section_divider(prs, 8, "Sign-off & Closure")

    # Acceptance Criteria Verification slide
    slide = create_content_slide(prs, "Acceptance Criteria Verification")
    verify_data = [
        ["#", "Criteria", "Target", "Actual", "Status"],
        ["1", "Log sources integrated", "[X]", "[X]", "PASS"],
        ["2", "Parse success rate", ">95%", "[X]%", "PASS"],
        ["3", "Use cases deployed", "[X]", "[X]", "PASS"],
        ["4", "False positive rate", "<20%", "[X]%", "PASS"],
        ["5", "Platform uptime", ">99.5%", "[X]%", "PASS"],
        ["6", "Training completed", "[X] staff", "[X] staff", "PASS"],
        ["7", "Documentation delivered", "[X] docs", "[X] docs", "PASS"],
    ]
    add_table(slide, verify_data, row_height=0.52)

    pass_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(5.8), Inches(11.8), Inches(0.4)
    )
    tf = pass_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "All acceptance criteria have been met."
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = COLORS["success_green"]
    run.font.name = "Segoe UI"

    # Project Sign-Off slide
    slide = create_content_slide(prs, "Project Sign-Off", "Acceptance Statement")

    statement_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(1.8), Inches(11.8), Inches(1.2)
    )
    tf = statement_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"By signing below, {customer_name} acknowledges that Seekurity has successfully completed the SIEM Implementation project in accordance with the Statement of Work dated [SOW_DATE], and all deliverables have been accepted."
    run.font.size = Pt(14)
    run.font.color.rgb = COLORS["dark_gray"]
    run.font.name = "Segoe UI"

    # Signature table
    sig_data = [
        ["", "Seekurity", customer_name],
        ["Name", "[NAME]", "[NAME]"],
        ["Title", "Project Manager", "[TITLE]"],
        ["Signature", "", ""],
        ["Date", "", ""],
    ]
    add_table(slide, sig_data, top=3.2, row_height=0.6)

    # =========================================================================
    # THANK YOU SLIDE
    # =========================================================================
    create_thank_you_slide(prs)

    # =========================================================================
    # APPENDIX
    # =========================================================================
    # Appendix A
    slide = create_content_slide(prs, "Appendix A: Document References")
    doc_data = [
        ["Document", "Location", "Version"],
        ["Statement of Work", "[LINK]", "[X]"],
        ["Solution Architecture", "[LINK]", "[X]"],
        ["Deployment Runbook", "[LINK]", "[X]"],
        ["Use Case Documentation", "[LINK]", "[X]"],
        ["Training Materials", "[LINK]", "[X]"],
        ["Change Log", "[LINK]", "[X]"],
    ]
    add_table(slide, doc_data, row_height=0.55)

    # Appendix B
    slide = create_content_slide(prs, "Appendix B: Log Source Inventory")
    inv_data = [
        ["#", "Source", "Type", "Method", "Status"],
        ["1", "[Source Name]", "[Type]", "[Method]", "Active"],
        ["2", "[Source Name]", "[Type]", "[Method]", "Active"],
        ["3", "[Source Name]", "[Type]", "[Method]", "Active"],
        ["...", "...", "...", "...", "..."],
    ]
    add_table(slide, inv_data, top=1.6, row_height=0.5)

    note_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.4), Inches(11.8), Inches(0.4)
    )
    tf = note_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Complete inventory available in Log Source Integration Guide"
    run.font.size = Pt(12)
    run.font.italic = True
    run.font.color.rgb = COLORS["light_gray"]
    run.font.name = "Segoe UI"

    # Appendix C
    slide = create_content_slide(prs, "Appendix C: Use Case Inventory")
    uc_data = [
        ["#", "Use Case", "Category", "Severity", "Status"],
        ["1", "[Use Case Name]", "[Category]", "[Sev]", "Active"],
        ["2", "[Use Case Name]", "[Category]", "[Sev]", "Active"],
        ["3", "[Use Case Name]", "[Category]", "[Sev]", "Active"],
        ["...", "...", "...", "...", "..."],
    ]
    add_table(slide, uc_data, top=1.6, row_height=0.5)

    note_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.4), Inches(11.8), Inches(0.4)
    )
    tf = note_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Complete inventory available in Use Case Documentation"
    run.font.size = Pt(12)
    run.font.italic = True
    run.font.color.rgb = COLORS["light_gray"]
    run.font.name = "Segoe UI"

    # Appendix D
    slide = create_content_slide(prs, "Appendix D: Version History")
    ver_data = [
        ["Version", "Date", "Author", "Changes"],
        ["0.1", "[DATE]", "[NAME]", "Initial draft"],
        ["0.2", "[DATE]", "[NAME]", "Added metrics"],
        ["1.0", "[DATE]", "[NAME]", "Final version"],
    ]
    add_table(slide, ver_data, row_height=0.55)

    # Save presentation
    prs.save(output_path)
    print(f"Completion presentation saved to: {output_path}")
    return prs


# =============================================================================
# MAIN EXECUTION
# =============================================================================

if __name__ == "__main__":
    # Get the directory of this script
    script_dir = os.path.dirname(os.path.abspath(__file__))

    # Generate presentations
    kickoff_path = os.path.join(
        script_dir, "SeekuritySIEM_Project_Kickoff_CUSTOMER_NAME.pptx"
    )
    completion_path = os.path.join(
        script_dir, "SeekuritySIEM_Project_Completion_CUSTOMER_NAME.pptx"
    )

    print("=" * 60)
    print("Seekurity SIEM Presentation Generator")
    print("=" * 60)
    print()

    generate_kickoff_presentation(kickoff_path)
    print()
    generate_completion_presentation(completion_path)

    print()
    print("=" * 60)
    print("Generation complete!")
    print("=" * 60)
