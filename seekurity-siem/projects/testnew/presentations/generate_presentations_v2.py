"""
Seekurity SIEM Professional Presentation Generator v2
Enhanced with SVG diagrams for enterprise-grade presentations.
"""

import io
import os
import tempfile

# Import diagram generators
from diagrams.architecture import (
    create_data_flow_diagram,
    create_network_topology_diagram,
    create_siem_architecture_background,
    create_siem_architecture_diagram,
    get_siem_architecture_text_positions,
)
from diagrams.metrics import (
    create_comparison_chart,
    create_donut_chart,
    create_kpi_dashboard,
    create_mini_kpi_card,
    create_pie_chart,
    create_progress_ring,
    create_risk_matrix,
    create_stat_list,
    create_status_summary,
    create_timeline_compact,
    create_vertical_bar_chart,
)
from diagrams.process import (
    create_escalation_pyramid,
    create_governance_model,
    create_implementation_phases,
    create_raci_matrix,
)
from diagrams.timeline import (
    create_gantt_chart,
    create_milestone_roadmap,
    create_phase_timeline,
)
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# =============================================================================
# BRAND COLORS (IBM-inspired)
# =============================================================================
COLORS = {
    "primary_blue": RGBColor(0x0F, 0x62, 0xFE),
    "secondary_blue": RGBColor(0x03, 0x53, 0xE9),
    "tertiary_blue": RGBColor(0x45, 0x89, 0xFF),
    "dark_gray": RGBColor(0x16, 0x16, 0x16),
    "medium_gray": RGBColor(0x52, 0x52, 0x52),
    "light_gray": RGBColor(0x6F, 0x6F, 0x6F),
    "background": RGBColor(0xFF, 0xFF, 0xFF),
    "alt_background": RGBColor(0xF4, 0xF4, 0xF4),
    "success_green": RGBColor(0x19, 0x80, 0x38),
    "warning_yellow": RGBColor(0xF1, 0xC2, 0x1B),
    "error_red": RGBColor(0xDA, 0x1E, 0x28),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "table_header": RGBColor(0x0F, 0x62, 0xFE),
    "table_row_alt": RGBColor(0xE8, 0xE8, 0xE8),
}

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def set_shape_fill(shape, color):
    """Set solid fill color for a shape."""
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


# Global Selenium driver (reused for performance)
_SELENIUM_DRIVER = None


def _get_selenium_driver():
    """Get or create a Selenium Chrome driver."""
    global _SELENIUM_DRIVER
    if _SELENIUM_DRIVER is None:
        try:
            from selenium import webdriver
            from selenium.webdriver.chrome.options import Options

            options = Options()
            options.add_argument("--headless")
            options.add_argument("--disable-gpu")
            options.add_argument("--no-sandbox")
            options.add_argument("--window-size=1600,1200")
            options.add_argument("--hide-scrollbars")

            _SELENIUM_DRIVER = webdriver.Chrome(options=options)
        except Exception as e:
            print(f"    Could not start Chrome: {e}")
            return None
    return _SELENIUM_DRIVER


def _cleanup_selenium():
    """Clean up Selenium driver."""
    global _SELENIUM_DRIVER
    if _SELENIUM_DRIVER:
        try:
            _SELENIUM_DRIVER.quit()
        except:
            pass
        _SELENIUM_DRIVER = None


def svg_to_png_selenium(svg_content: str, output_path: str) -> bool:
    """
    Convert SVG to PNG using Selenium with headless Chrome.
    This method works on Windows without Cairo.
    """
    driver = _get_selenium_driver()
    if not driver:
        return False

    html_path = None
    try:
        from selenium.webdriver.common.by import By

        # Create HTML with SVG
        html = f"""<!DOCTYPE html>
        <html><head>
        <style>
            body {{ margin: 0; padding: 10px; background: white; }}
            svg {{ display: block; }}
        </style>
        </head>
        <body>{svg_content}</body></html>"""

        # Save to temp file
        with tempfile.NamedTemporaryFile(
            mode="w", suffix=".html", delete=False, encoding="utf-8"
        ) as f:
            f.write(html)
            html_path = f.name

        # Load HTML
        driver.get("file:///" + html_path.replace("\\", "/"))

        # Find SVG element and screenshot it
        svg_elem = driver.find_element(By.TAG_NAME, "svg")
        svg_elem.screenshot(output_path)

        return os.path.exists(output_path)

    except Exception as e:
        # Silently fail
        pass
    finally:
        if html_path and os.path.exists(html_path):
            try:
                os.unlink(html_path)
            except:
                pass

    return False


def add_svg_to_slide(
    slide,
    svg_content: str,
    left: float,
    top: float,
    width: float = None,
    height: float = None,
    diagram_name: str = None,
):
    """
    Add an SVG diagram to a slide as a PNG image.
    Uses Selenium with headless Chrome for conversion.
    """
    temp_png = tempfile.mktemp(suffix=".png")

    try:
        # Convert SVG to PNG using Selenium
        success = svg_to_png_selenium(svg_content, temp_png)

        if success and os.path.exists(temp_png):
            png_size = os.path.getsize(temp_png)
            print(
                f"    + Diagram embedded: {diagram_name or 'unnamed'} ({png_size} bytes)"
            )
            # Add image to slide
            if width and height:
                slide.shapes.add_picture(
                    temp_png, Inches(left), Inches(top), Inches(width), Inches(height)
                )
            elif width:
                slide.shapes.add_picture(
                    temp_png, Inches(left), Inches(top), width=Inches(width)
                )
            else:
                slide.shapes.add_picture(temp_png, Inches(left), Inches(top))
            return True
        else:
            print(
                f"    ! Diagram failed: {diagram_name or 'unnamed'} (success={success})"
            )

    except Exception as e:
        print(f"    ! Diagram error: {diagram_name or 'unnamed'} - {e}")
    finally:
        # Clean up temp file
        if os.path.exists(temp_png):
            try:
                os.unlink(temp_png)
            except:
                pass

    # Fallback: Add placeholder rectangle
    placeholder_width = width or 10
    placeholder_height = height or 4

    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(placeholder_width),
        Inches(placeholder_height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF4, 0xF4, 0xF4)
    shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    shape.line.width = Pt(1)

    # Add text to placeholder
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "[Diagram]"
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(0x6F, 0x6F, 0x6F)
    run.font.name = "Segoe UI"

    return False


def export_svg_diagrams(output_dir: str, diagrams: dict):
    """Export SVG diagrams to files."""
    if not diagrams:
        return

    svg_dir = os.path.join(output_dir, "diagrams_svg")
    os.makedirs(svg_dir, exist_ok=True)

    for name, content in SVG_DIAGRAMS.items():
        save_svg_file(content, name, svg_dir)

    print(f"  Exported {len(SVG_DIAGRAMS)} SVG diagrams to: {svg_dir}")


def create_title_slide(prs, title, subtitle, customer_name, date, company="Seekurity"):
    """Create a title slide with blue background."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Blue background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    set_shape_fill(background, COLORS["primary_blue"])
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(2.2), Inches(11.8), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = COLORS["white"]
    run.font.name = "Segoe UI"

    # Subtitle
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(12)
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.size = Pt(28)
    run2.font.color.rgb = COLORS["white"]
    run2.font.name = "Segoe UI"

    # Customer name
    customer_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(4.2), Inches(11.8), Inches(0.6)
    )
    tf2 = customer_box.text_frame
    p3 = tf2.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    run3 = p3.add_run()
    run3.text = customer_name
    run3.font.size = Pt(24)
    run3.font.bold = True
    run3.font.color.rgb = COLORS["white"]
    run3.font.name = "Segoe UI"

    # Company and date
    footer_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(5.5), Inches(11.8), Inches(0.8)
    )
    tf3 = footer_box.text_frame
    p4 = tf3.paragraphs[0]
    p4.alignment = PP_ALIGN.CENTER
    run4 = p4.add_run()
    run4.text = f"Prepared by {company}"
    run4.font.size = Pt(16)
    run4.font.color.rgb = COLORS["white"]
    run4.font.name = "Segoe UI"

    p5 = tf3.add_paragraph()
    p5.alignment = PP_ALIGN.CENTER
    p5.space_before = Pt(6)
    run5 = p5.add_run()
    run5.text = date
    run5.font.size = Pt(14)
    run5.font.color.rgb = COLORS["white"]
    run5.font.name = "Segoe UI"

    return slide


def create_section_divider(prs, section_number, section_title):
    """Create a section divider slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Gray background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    set_shape_fill(background, COLORS["alt_background"])
    background.line.fill.background()

    # Section number
    num_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(2.5), Inches(11.8), Inches(1)
    )
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{section_number:02d}"
    run.font.size = Pt(72)
    run.font.bold = True
    run.font.color.rgb = COLORS["primary_blue"]
    run.font.name = "Segoe UI"

    # Section title
    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(3.8), Inches(11.8), Inches(0.8)
    )
    tf2 = title_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = section_title
    run2.font.size = Pt(36)
    run2.font.color.rgb = COLORS["dark_gray"]
    run2.font.name = "Segoe UI"

    return slide


def create_content_slide(prs, title, subtitle=None):
    """Create a content slide with title."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(0.5), Inches(11.8), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = COLORS["primary_blue"]
    run.font.name = "Segoe UI"

    if subtitle:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(6)
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(20)
        run2.font.color.rgb = COLORS["medium_gray"]
        run2.font.name = "Segoe UI"

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(1.35), Inches(1.5), Inches(0.04)
    )
    set_shape_fill(line, COLORS["primary_blue"])
    line.line.fill.background()

    return slide


def create_diagram_slide(prs, title, svg_content, subtitle=None):
    """Create a slide with a diagram."""
    slide = create_content_slide(prs, title, subtitle)
    add_svg_to_slide(slide, svg_content, left=0.5, top=1.6, width=12.3)
    return slide


def add_table(
    slide,
    data,
    left=0.75,
    top=1.6,
    width=11.8,
    row_height=0.45,
    header_color=None,
    font_size=12,
):
    """Add a formatted table to a slide."""
    rows = len(data)
    cols = len(data[0]) if data else 0

    if rows == 0 or cols == 0:
        return None

    col_width = width / cols
    table_height = rows * row_height

    table = slide.shapes.add_table(
        rows, cols, Inches(left), Inches(top), Inches(width), Inches(table_height)
    ).table

    for i, row_data in enumerate(data):
        for j, cell_text in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = str(cell_text)

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(font_size)
                paragraph.font.name = "Segoe UI"

                if i == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = COLORS["white"]
                else:
                    paragraph.font.color.rgb = COLORS["dark_gray"]

            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_color or COLORS["table_header"]
            elif i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["table_row_alt"]
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["white"]

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table


def create_thank_you_slide(
    prs,
    company="Seekurity",
    tagline="Enterprise Security Solutions",
    email="contact@seekurity.com",
    phone="+1 (555) 123-4567",
    website="www.seekurity.com",
):
    """Create a thank you slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Blue background
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )
    set_shape_fill(background, COLORS["primary_blue"])
    background.line.fill.background()

    # Thank you text
    title_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(2.2), Inches(11.8), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Thank You"
    run.font.size = Pt(56)
    run.font.bold = True
    run.font.color.rgb = COLORS["white"]
    run.font.name = "Segoe UI"

    # Company
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(24)
    run2 = p2.add_run()
    run2.text = company
    run2.font.size = Pt(28)
    run2.font.bold = True
    run2.font.color.rgb = COLORS["white"]
    run2.font.name = "Segoe UI"

    # Tagline
    p3 = tf.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(6)
    run3 = p3.add_run()
    run3.text = tagline
    run3.font.size = Pt(18)
    run3.font.color.rgb = COLORS["white"]
    run3.font.name = "Segoe UI"

    # Contact
    contact_box = slide.shapes.add_textbox(
        Inches(0.75), Inches(5.2), Inches(11.8), Inches(1.2)
    )
    tf2 = contact_box.text_frame
    for text in [email, phone, website]:
        p = tf2.add_paragraph() if tf2.paragraphs[0].text else tf2.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = text
        run.font.size = Pt(14)
        run.font.color.rgb = COLORS["white"]
        run.font.name = "Segoe UI"

    return slide


# =============================================================================
# PROFESSIONAL LAYOUT FUNCTIONS (Magazine Style)
# =============================================================================


def create_header_bar(slide, left_text, right_text):
    """
    Create a header bar with breadcrumb-style navigation.
    Like: [icon] LEFT TEXT                    RIGHT TEXT
    """
    # Header background bar
    header = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.4),
        Inches(0.25),
        Inches(12.5),
        Inches(0.55),
    )
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS["white"]
    header.line.color.rgb = COLORS["primary_blue"]
    header.line.width = Pt(2)

    # Left text (with home icon placeholder)
    left_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(0.35), Inches(6), Inches(0.4)
    )
    tf = left_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"⌂  {left_text}"
    run.font.size = Pt(11)
    run.font.color.rgb = COLORS["dark_gray"]
    run.font.name = "Segoe UI"

    # Right text
    right_box = slide.shapes.add_textbox(
        Inches(9), Inches(0.35), Inches(3.5), Inches(0.4)
    )
    tf2 = right_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = right_text
    run2.font.size = Pt(11)
    run2.font.bold = True
    run2.font.color.rgb = COLORS["dark_gray"]
    run2.font.name = "Segoe UI"


def create_magazine_slide_type1(
    prs,
    header_left,
    header_right,
    hero_title,
    section_title,
    body_text,
    methodology_title,
    methodology_text,
):
    """
    Magazine-style layout Type 1: Research Approach style
    - Header bar
    - Large hero title (left)
    - Image placeholder (top right)
    - Section with title and body (left)
    - Colored section (bottom right) with light background diagonal
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Light background with diagonal accent
    bg_accent = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM, Inches(6), Inches(4), Inches(8), Inches(4)
    )
    bg_accent.fill.solid()
    bg_accent.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xE8)  # Light blue-ish tint
    bg_accent.line.fill.background()
    bg_accent.rotation = -5

    # Header bar
    create_header_bar(slide, header_left, header_right)

    # Hero title (large, left side)
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.1), Inches(7), Inches(1.2)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = hero_title
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = COLORS["primary_blue"]
    run.font.name = "Segoe UI"

    # Image placeholder (top right) - rounded rectangle
    img_placeholder = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8), Inches(1.1), Inches(4.8), Inches(2.4)
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = RGBColor(0xE0, 0xE8, 0xF0)
    img_placeholder.line.color.rgb = RGBColor(0xC0, 0xD0, 0xE0)
    img_placeholder.line.width = Pt(1)

    # Image placeholder text
    img_text = slide.shapes.add_textbox(
        Inches(8), Inches(2.0), Inches(4.8), Inches(0.5)
    )
    tf = img_text.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "[Image Placeholder]"
    run.font.size = Pt(12)
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    # Section title (left, below hero)
    section_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(6.5), Inches(0.4)
    )
    tf = section_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = section_title
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = COLORS["primary_blue"]
    run.font.name = "Segoe UI"

    # Body text
    body_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.0), Inches(6.5), Inches(1.5)
    )
    tf = body_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = body_text
    run.font.size = Pt(12)
    run.font.color.rgb = COLORS["dark_gray"]
    run.font.name = "Segoe UI"

    # Second image placeholder (bottom left)
    img2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.5), Inches(5.5), Inches(2.5)
    )
    img2.fill.solid()
    img2.fill.fore_color.rgb = RGBColor(0xE0, 0xE8, 0xF0)
    img2.line.color.rgb = RGBColor(0xC0, 0xD0, 0xE0)
    img2.line.width = Pt(1)

    # Methodology section (bottom right)
    method_title_box = slide.shapes.add_textbox(
        Inches(7), Inches(4.8), Inches(5.5), Inches(0.4)
    )
    tf = method_title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = methodology_title
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = COLORS["dark_gray"]
    run.font.name = "Segoe UI"

    method_body = slide.shapes.add_textbox(
        Inches(7), Inches(5.3), Inches(5.5), Inches(1.8)
    )
    tf = method_body.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = methodology_text
    run.font.size = Pt(11)
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    return slide


def create_magazine_slide_type2(prs, header_left, header_right, hero_title, cards):
    """
    Magazine-style layout Type 2: Cards grid with hero title
    - Header bar
    - Large hero title (left, 2 lines)
    - Image placeholder (right)
    - 2x2 grid of cards (2 white, 2 colored)

    Args:
        cards: List of 4 dicts with 'title', 'body', 'colored' (bool)
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Background - light tint on right side
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(8.5),
        Inches(0),
        SLIDE_WIDTH - Inches(8.5),
        SLIDE_HEIGHT,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS["primary_blue"]
    bg.fill.fore_color.brightness = 0.4
    bg.line.fill.background()

    # Header bar
    create_header_bar(slide, header_left, header_right)

    # Hero title (large, uppercase, 2 lines)
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.1), Inches(7), Inches(1.6)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = hero_title.upper()
    run.font.size = Pt(42)
    run.font.bold = True
    run.font.color.rgb = COLORS["dark_gray"]
    run.font.name = "Segoe UI"

    # Image placeholder (right side)
    img = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.2), Inches(1.3), Inches(4.5), Inches(3.5)
    )
    img.fill.solid()
    img.fill.fore_color.rgb = RGBColor(0xD8, 0xE8, 0xF8)
    img.line.fill.background()

    # 2x2 Cards grid
    card_positions = [
        (0.5, 3.0, False),  # Top left - white
        (4.2, 3.0, False),  # Top right - white
        (0.5, 5.0, True),  # Bottom left - colored
        (4.2, 5.0, True),  # Bottom right - colored
    ]

    for i, (left, top, colored) in enumerate(card_positions):
        if i < len(cards):
            card = cards[i]
            _add_magazine_card(
                slide,
                card.get("title", ""),
                card.get("body", ""),
                left,
                top,
                width=3.5,
                height=1.8,
                colored=card.get("colored", colored),
            )

    return slide


def _add_magazine_card(
    slide, title, body, left, top, width=3.5, height=1.8, colored=False
):
    """Add a magazine-style content card."""
    # Card background
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )

    if colored:
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS["primary_blue"]
        card.line.fill.background()
        title_color = COLORS["white"]
        body_color = RGBColor(0xE0, 0xE8, 0xFF)
    else:
        card.fill.solid()
        card.fill.fore_color.rgb = COLORS["white"]
        card.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        card.line.width = Pt(1)
        title_color = COLORS["primary_blue"]
        body_color = COLORS["dark_gray"]

    # Card title
    title_box = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.15), Inches(width - 0.3), Inches(0.4)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = title_color
    run.font.name = "Segoe UI"

    # Card body
    body_box = slide.shapes.add_textbox(
        Inches(left + 0.15),
        Inches(top + 0.55),
        Inches(width - 0.3),
        Inches(height - 0.7),
    )
    tf = body_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = body
    run.font.size = Pt(10)
    run.font.color.rgb = body_color
    run.font.name = "Segoe UI"


def create_magazine_slide_type3(
    prs, header_left, header_right, quote_text, section_title, bullet_points
):
    """
    Magazine-style layout Type 3: Quote + Bullet points
    - Header bar
    - Large quote text (top, colored)
    - Bottom section with gray background
    - Section title + bullet points
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Bottom gray section
    bottom_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(4.2), SLIDE_WIDTH, Inches(3.3)
    )
    bottom_bg.fill.solid()
    bottom_bg.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xF0)
    bottom_bg.line.fill.background()

    # Header bar
    create_header_bar(slide, header_left, header_right)

    # Quote text (large, colored, italic-style)
    quote_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2), Inches(12), Inches(2.5)
    )
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = quote_text.upper()
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = COLORS["primary_blue"]
    run.font.name = "Segoe UI"

    # Section title (in gray area)
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.5), Inches(5), Inches(0.8)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = section_title
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = COLORS["dark_gray"]
    run.font.name = "Segoe UI"

    # Bullet points (right side of gray area)
    bullets_box = slide.shapes.add_textbox(
        Inches(5.5), Inches(4.5), Inches(7), Inches(2.5)
    )
    tf = bullets_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(bullet_points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        run = p.add_run()
        run.text = f"• {point}"
        run.font.size = Pt(12)
        run.font.color.rgb = COLORS["dark_gray"]
        run.font.name = "Segoe UI"

    return slide


def create_magazine_slide_type4(
    prs, header_left, header_right, hero_title, subtitle, body_text, side_cards
):
    """
    Magazine-style layout Type 4: Emerging Technologies style
    - Header bar
    - Large hero title (left)
    - Image placeholder (center)
    - Stacked cards on right (3 cards)
    - Subtitle + body text (bottom left)

    Args:
        side_cards: List of 3 dicts with 'title', 'body'
    """
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Right side colored background
    right_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(8.3), Inches(0), Inches(5), SLIDE_HEIGHT
    )
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = COLORS["primary_blue"]
    right_bg.fill.fore_color.brightness = 0.35
    right_bg.line.fill.background()

    # Header bar
    create_header_bar(slide, header_left, header_right)

    # Hero title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.2), Inches(5), Inches(1.6)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = hero_title
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = COLORS["primary_blue"]
    run.font.name = "Segoe UI"

    # Center image placeholder
    img = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.5), Inches(1.2), Inches(3.5), Inches(4.5)
    )
    img.fill.solid()
    img.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xF8)
    img.line.color.rgb = RGBColor(0xC0, 0xD8, 0xE8)
    img.line.width = Pt(2)

    # Subtitle (bottom left)
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(4), Inches(0.4))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = subtitle
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = COLORS["primary_blue"]
    run.font.name = "Segoe UI"

    # Body text
    body_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(4), Inches(2))
    tf = body_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = body_text
    run.font.size = Pt(11)
    run.font.color.rgb = COLORS["dark_gray"]
    run.font.name = "Segoe UI"

    # Right side stacked cards
    card_tops = [1.3, 3.1, 5.0]
    for i, top in enumerate(card_tops):
        if i < len(side_cards):
            card = side_cards[i]
            _add_stacked_card(
                slide,
                card.get("title", f"Item {i + 1}"),
                card.get("body", ""),
                left=8.5,
                top=top,
                width=4.3,
                height=1.6,
            )

    return slide


def _add_stacked_card(slide, title, body, left, top, width, height):
    """Add a stacked card for Type 4 layout."""
    # Card background (slightly lighter blue)
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["primary_blue"]
    card.fill.fore_color.brightness = 0.15
    card.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.12), Inches(width - 0.3), Inches(0.35)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = COLORS["white"]
    run.font.name = "Segoe UI"

    # Body
    body_box = slide.shapes.add_textbox(
        Inches(left + 0.15),
        Inches(top + 0.5),
        Inches(width - 0.3),
        Inches(height - 0.6),
    )
    tf = body_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = body
    run.font.size = Pt(10)
    run.font.color.rgb = RGBColor(0xE8, 0xF0, 0xFF)
    run.font.name = "Segoe UI"


# =============================================================================
# COMPLEX LAYOUT FUNCTIONS
# =============================================================================


def create_split_slide_50_50(prs, title, subtitle=None):
    """Create a slide with 50/50 split layout (left and right panels)."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = COLORS["primary_blue"]
    run.font.name = "Segoe UI"

    if subtitle:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(16)
        run2.font.color.rgb = COLORS["medium_gray"]
        run2.font.name = "Segoe UI"

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(1.2), Inches(0.035)
    )
    set_shape_fill(line, COLORS["primary_blue"])
    line.line.fill.background()

    # Left panel background (subtle)
    left_panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.3), Inches(6.1), Inches(5.9)
    )
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
    left_panel.line.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    left_panel.line.width = Pt(1)

    # Right panel background
    right_panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.3), Inches(6.1), Inches(5.9)
    )
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
    right_panel.line.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    right_panel.line.width = Pt(1)

    return slide


def add_info_card(slide, title, items, left, top, width=5.8, card_color=None):
    """Add an information card with title and bullet items."""
    card_height = 0.5 + len(items) * 0.4 + 0.3

    # Card background
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(card_height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["white"]
    card.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    card.line.width = Pt(1)

    # Card title bar
    title_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(0.45),
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = card_color or COLORS["primary_blue"]
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(
        Inches(left + 0.15), Inches(top + 0.08), Inches(width - 0.3), Inches(0.35)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = COLORS["white"]
    run.font.name = "Segoe UI"

    # Items
    items_box = slide.shapes.add_textbox(
        Inches(left + 0.15),
        Inches(top + 0.5),
        Inches(width - 0.3),
        Inches(card_height - 0.55),
    )
    tf = items_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        p.space_after = Pt(3)
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(11)
        run.font.color.rgb = COLORS["dark_gray"]
        run.font.name = "Segoe UI"

    return card_height


def add_metric_box(slide, value, label, left, top, width=2.8, height=1.2, color=None):
    """Add a metric display box with large value and label."""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    box.fill.solid()
    box.fill.fore_color.rgb = color or COLORS["primary_blue"]
    box.line.fill.background()

    # Value
    value_box = slide.shapes.add_textbox(
        Inches(left), Inches(top + 0.15), Inches(width), Inches(0.6)
    )
    tf = value_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(value)
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = COLORS["white"]
    run.font.name = "Segoe UI"

    # Label
    label_box = slide.shapes.add_textbox(
        Inches(left), Inches(top + 0.7), Inches(width), Inches(0.4)
    )
    tf2 = label_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(11)
    run2.font.color.rgb = COLORS["white"]
    run2.font.name = "Segoe UI"


def add_detail_table(slide, headers, rows, left, top, width=5.8, font_size=10):
    """Add a compact detail table."""
    cols = len(headers)
    total_rows = len(rows) + 1
    row_height = 0.35

    table = slide.shapes.add_table(
        total_rows,
        cols,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(total_rows * row_height),
    ).table

    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLORS["primary_blue"]
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = COLORS["white"]
            p.font.name = "Segoe UI"
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                COLORS["white"] if i % 2 == 0 else COLORS["table_row_alt"]
            )
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = COLORS["dark_gray"]
                p.font.name = "Segoe UI"
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table


def add_status_indicator(slide, status, label, left, top, width=2.5):
    """Add a status indicator with colored dot and label."""
    colors = {
        "complete": COLORS["success_green"],
        "in_progress": COLORS["primary_blue"],
        "pending": COLORS["warning_yellow"],
        "blocked": COLORS["error_red"],
    }

    # Dot
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left),
        Inches(top + 0.05),
        Inches(0.15),
        Inches(0.15),
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = colors.get(status, COLORS["medium_gray"])
    dot.line.fill.background()

    # Label
    label_box = slide.shapes.add_textbox(
        Inches(left + 0.22), Inches(top), Inches(width - 0.22), Inches(0.25)
    )
    tf = label_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label
    run.font.size = Pt(11)
    run.font.color.rgb = COLORS["dark_gray"]
    run.font.name = "Segoe UI"


def create_split_slide_60_40(prs, title, subtitle=None):
    """Create a slide with 60/40 split layout (larger left, smaller right)."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = COLORS["primary_blue"]
    run.font.name = "Segoe UI"

    if subtitle:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(16)
        run2.font.color.rgb = COLORS["medium_gray"]
        run2.font.name = "Segoe UI"

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(1.2), Inches(0.035)
    )
    set_shape_fill(line, COLORS["primary_blue"])
    line.line.fill.background()

    # Left panel (60%)
    left_panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.3), Inches(7.4), Inches(5.9)
    )
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
    left_panel.line.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    left_panel.line.width = Pt(1)

    # Right panel (40%)
    right_panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(8.0), Inches(1.3), Inches(4.9), Inches(5.9)
    )
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
    right_panel.line.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
    right_panel.line.width = Pt(1)

    return slide


def create_dashboard_slide(prs, title, subtitle=None):
    """Create a slide with dashboard layout (metrics row + content area)."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = COLORS["primary_blue"]
    run.font.name = "Segoe UI"

    if subtitle:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(16)
        run2.font.color.rgb = COLORS["medium_gray"]
        run2.font.name = "Segoe UI"

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(1.2), Inches(0.035)
    )
    set_shape_fill(line, COLORS["primary_blue"])
    line.line.fill.background()

    return slide


def create_three_column_slide(prs, title, subtitle=None):
    """Create a slide with three equal columns."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = COLORS["primary_blue"]
    run.font.name = "Segoe UI"

    if subtitle:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(4)
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(16)
        run2.font.color.rgb = COLORS["medium_gray"]
        run2.font.name = "Segoe UI"

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(1.2), Inches(0.035)
    )
    set_shape_fill(line, COLORS["primary_blue"])
    line.line.fill.background()

    # Three column backgrounds
    for i in range(3):
        col_panel = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.4 + i * 4.2),
            Inches(1.3),
            Inches(3.95),
            Inches(5.9),
        )
        col_panel.fill.solid()
        col_panel.fill.fore_color.rgb = RGBColor(0xF8, 0xF8, 0xF8)
        col_panel.line.color.rgb = RGBColor(0xE8, 0xE8, 0xE8)
        col_panel.line.width = Pt(1)

    return slide


def add_phase_card(slide, phase_num, title, status, items, left, top, width=3.8):
    """Add a phase card with number, title, status, and details."""
    card_height = 2.2 + len(items) * 0.28

    # Card background
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(card_height),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = COLORS["white"]
    card.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    card.line.width = Pt(1)

    # Phase number circle
    status_colors = {
        "complete": COLORS["success_green"],
        "in_progress": COLORS["primary_blue"],
        "pending": COLORS["light_gray"],
    }
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left + width / 2 - 0.35),
        Inches(top + 0.15),
        Inches(0.7),
        Inches(0.7),
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = status_colors.get(status, COLORS["light_gray"])
    circle.line.fill.background()

    # Phase number
    num_box = slide.shapes.add_textbox(
        Inches(left + width / 2 - 0.35), Inches(top + 0.28), Inches(0.7), Inches(0.5)
    )
    tf = num_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(phase_num)
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = COLORS["white"]
    run.font.name = "Segoe UI"

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 0.95), Inches(width - 0.2), Inches(0.4)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = COLORS["dark_gray"]
    run.font.name = "Segoe UI"

    # Status label
    status_labels = {
        "complete": "Completed",
        "in_progress": "In Progress",
        "pending": "Pending",
    }
    status_box = slide.shapes.add_textbox(
        Inches(left + 0.1), Inches(top + 1.3), Inches(width - 0.2), Inches(0.3)
    )
    tf = status_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = status_labels.get(status, status)
    run.font.size = Pt(10)
    run.font.color.rgb = status_colors.get(status, COLORS["medium_gray"])
    run.font.name = "Segoe UI"

    # Divider
    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left + 0.3),
        Inches(top + 1.65),
        Inches(width - 0.6),
        Inches(0.015),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    divider.line.fill.background()

    # Items
    items_box = slide.shapes.add_textbox(
        Inches(left + 0.15),
        Inches(top + 1.75),
        Inches(width - 0.3),
        Inches(card_height - 1.9),
    )
    tf = items_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = f"• {item}"
        run.font.size = Pt(9)
        run.font.color.rgb = COLORS["medium_gray"]
        run.font.name = "Segoe UI"


def add_key_value_list(slide, items, left, top, width=5.5, label_width=2.0):
    """Add a list of key-value pairs."""
    items_box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(len(items) * 0.35)
    )
    tf = items_box.text_frame
    tf.word_wrap = True

    for i, (key, value) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(4)

        run = p.add_run()
        run.text = f"{key}: "
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = COLORS["dark_gray"]
        run.font.name = "Segoe UI"

        run2 = p.add_run()
        run2.text = value
        run2.font.size = Pt(11)
        run2.font.color.rgb = COLORS["medium_gray"]
        run2.font.name = "Segoe UI"


# =============================================================================
# EDITABLE DIAGRAM FUNCTIONS
# =============================================================================


def create_editable_architecture_slide(prs, title, subtitle=None):
    """
    Create a slide with the SIEM architecture diagram using editable text boxes.
    The diagram background is a PNG image, with text as editable PowerPoint text boxes.

    Args:
        prs: Presentation object
        title: Slide title
        subtitle: Optional subtitle

    Returns:
        slide: The created slide
    """
    slide = create_content_slide(prs, title, subtitle)

    # Get background SVG (no text)
    arch_bg = create_siem_architecture_background(width=1100, height=480)

    # Add background image
    add_svg_to_slide(
        slide, arch_bg, left=0.4, top=1.5, width=12.5, diagram_name="arch_background"
    )

    # Add editable text boxes on top
    # Scale factor for 12.5" wide display of 1100px diagram
    s = 12.5 / 1100
    base_left = 0.4
    base_top = 1.5

    # Layer labels (gray, outside boxes)
    _add_layer_label(slide, "PRESENTATION TIER", base_left + s * 20, base_top + s * 5)
    _add_layer_label(slide, "PROCESSING TIER", base_left + s * 20, base_top + s * 125)
    _add_layer_label(slide, "COLLECTION TIER", base_left + s * 20, base_top + s * 245)
    _add_layer_label(slide, "DATA SOURCES", base_left + s * 20, base_top + s * 365)

    # Presentation tier boxes (white text on blue)
    _add_box_label(
        slide,
        "Dashboards",
        "Real-time Views",
        base_left + s * 80,
        base_top + s * 45,
        s * 160,
        s * 55,
    )
    _add_box_label(
        slide,
        "Alerts",
        "Notifications",
        base_left + s * 270,
        base_top + s * 45,
        s * 160,
        s * 55,
    )
    _add_box_label(
        slide,
        "Reports",
        "Scheduled",
        base_left + s * 460,
        base_top + s * 45,
        s * 160,
        s * 55,
    )
    _add_box_label(
        slide,
        "Investigation",
        "Ad-hoc Search",
        base_left + s * 650,
        base_top + s * 45,
        s * 160,
        s * 55,
    )

    # Processing tier boxes
    _add_box_label(
        slide,
        "Indexer Cluster",
        "Index & Store Events",
        base_left + s * 100,
        base_top + s * 165,
        s * 200,
        s * 55,
    )
    _add_box_label(
        slide,
        "Correlation Engine",
        "Rules & Analytics",
        base_left + s * 350,
        base_top + s * 165,
        s * 200,
        s * 55,
    )
    _add_box_label(
        slide,
        "Search Cluster",
        "Distributed Search",
        base_left + s * 600,
        base_top + s * 165,
        s * 200,
        s * 55,
    )

    # Collection tier boxes
    _add_box_label(
        slide,
        "Log Collectors",
        "Syslog / Agent / API",
        base_left + s * 140,
        base_top + s * 285,
        s * 180,
        s * 55,
    )
    _add_box_label(
        slide,
        "Heavy Forwarders",
        "Parse & Transform",
        base_left + s * 360,
        base_top + s * 285,
        s * 180,
        s * 55,
    )
    _add_box_label(
        slide,
        "Load Balancer",
        "HA Distribution",
        base_left + s * 580,
        base_top + s * 285,
        s * 180,
        s * 55,
    )

    # Data source boxes
    _add_box_label(
        slide,
        "Firewall",
        "Palo Alto",
        base_left + s * 50,
        base_top + s * 410,
        s * 100,
        s * 50,
        size=11,
    )
    _add_box_label(
        slide,
        "EDR",
        "CrowdStrike",
        base_left + s * 170,
        base_top + s * 410,
        s * 100,
        s * 50,
        size=11,
    )
    _add_box_label(
        slide,
        "Identity",
        "Azure AD",
        base_left + s * 290,
        base_top + s * 410,
        s * 100,
        s * 50,
        size=11,
    )
    _add_box_label(
        slide,
        "Cloud",
        "AWS/Azure",
        base_left + s * 410,
        base_top + s * 410,
        s * 100,
        s * 50,
        size=11,
    )
    _add_box_label(
        slide,
        "Email",
        "O365",
        base_left + s * 530,
        base_top + s * 410,
        s * 100,
        s * 50,
        size=11,
    )
    _add_box_label(
        slide,
        "Apps",
        "Custom",
        base_left + s * 650,
        base_top + s * 410,
        s * 100,
        s * 50,
        size=11,
    )
    _add_box_label(
        slide,
        "+ More",
        None,
        base_left + s * 770,
        base_top + s * 410,
        s * 80,
        s * 50,
        size=11,
        color="gray",
    )

    # Data flow annotation
    _add_layer_label(
        slide, "Data Flow", base_left + s * 950, base_top + s * 180, size=10
    )
    _add_layer_label(slide, "[X] EPS", base_left + s * 950, base_top + s * 280, size=10)

    return slide


def _add_layer_label(slide, text, left, top, size=10):
    """Add a layer label (gray text, uppercase style)."""
    text_box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(1.5), Inches(0.25)
    )
    tf = text_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = COLORS["medium_gray"]
    run.font.name = "Segoe UI"


def _add_box_label(
    slide, text, subtitle, left, top, width, height, size=12, color="white"
):
    """Add a box label with optional subtitle (centered, for diagram boxes)."""
    text_box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = text_box.text_frame
    tf.word_wrap = True

    # Center vertically by adding space at top
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(8 if subtitle else 12)

    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = COLORS["white"] if color == "white" else COLORS["medium_gray"]
    run.font.name = "Segoe UI"

    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(2)
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.size = Pt(size - 2)
        run2.font.color.rgb = (
            RGBColor(0xE0, 0xE0, 0xE0) if color == "white" else COLORS["light_gray"]
        )
        run2.font.name = "Segoe UI"


# =============================================================================
# KICKOFF PRESENTATION GENERATOR
# =============================================================================


def generate_kickoff_presentation(
    output_path, customer_name="[CUSTOMER_NAME]", date="[DATE]"
):
    """Generate the Project Kickoff presentation with diagrams."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    print("Generating Kickoff Presentation...")

    # Title Slide
    create_title_slide(
        prs, "SIEM Implementation", "Project Kickoff", customer_name, date
    )

    # Agenda
    slide = create_content_slide(prs, "Agenda")
    agenda_data = [
        ["#", "Topic", "Duration"],
        ["1", "Introduction & Objectives", "10 min"],
        ["2", "Project Scope & Deliverables", "15 min"],
        ["3", "Technical Architecture", "20 min"],
        ["4", "Implementation Approach", "15 min"],
        ["5", "Timeline & Milestones", "10 min"],
        ["6", "Governance & Communication", "10 min"],
        ["7", "Risks & Mitigation", "10 min"],
        ["8", "Next Steps & Q&A", "10 min"],
    ]
    add_table(slide, agenda_data)

    # Section 1: Introduction
    create_section_divider(prs, 1, "Introduction & Objectives")

    # NEW: Magazine-style Type 1 - Project Introduction
    create_magazine_slide_type1(
        prs,
        header_left="SIEM IMPLEMENTATION PROJECT",
        header_right="PROJECT KICKOFF",
        hero_title="Project\nIntroduction",
        section_title="Background & Context",
        body_text="This SIEM implementation project aims to establish a comprehensive security monitoring capability. The solution will integrate log sources across your enterprise infrastructure, enabling real-time threat detection and rapid incident response.",
        methodology_title="Our Approach",
        methodology_text="We follow a proven methodology combining industry best practices with tailored solutions. Our team will work closely with your security and IT teams to ensure seamless integration and knowledge transfer throughout the project lifecycle.",
    )

    slide = create_content_slide(prs, "Meeting Objectives")
    add_table(
        slide,
        [
            ["Objective", "Outcome"],
            ["Alignment", "Confirm project goals and success criteria"],
            ["Scope", "Review deliverables and exclusions"],
            ["Architecture", "Understand technical design and integrations"],
            ["Timeline", "Agree on milestones and dependencies"],
            ["Governance", "Establish communication and escalation paths"],
            ["Risks", "Identify and plan mitigations"],
        ],
    )

    # Section 2: Scope
    create_section_divider(prs, 2, "Project Scope & Deliverables")

    slide = create_content_slide(prs, "Scope Overview")
    add_table(
        slide,
        [
            ["Category", "In Scope", "Details"],
            ["Platform", "Yes", "SIEM deployment and configuration"],
            ["Log Sources", "Yes", "[X] sources across [Y] categories"],
            ["Use Cases", "Yes", "[Z] detection rules and alerts"],
            ["Integrations", "Yes", "Ticketing, SOAR, threat intelligence"],
            ["Training", "Yes", "Administrator and analyst sessions"],
            ["Managed Services", "No", "Post-warranty support"],
        ],
    )

    # NEW: Project Overview - Split Layout 50/50
    slide = create_split_slide_50_50(
        prs, "Project Overview", "Scope and Objectives at a Glance"
    )

    # Left side: Project details card
    add_info_card(
        slide,
        "Project Details",
        [
            f"Customer: {customer_name}",
            "Project Type: SIEM Implementation",
            "Duration: 12 weeks",
            "Team Size: 8 resources",
            "Methodology: Agile/Iterative",
        ],
        left=0.55,
        top=1.5,
        width=5.8,
        card_color=COLORS["primary_blue"],
    )

    # Left side: Success criteria
    add_info_card(
        slide,
        "Success Criteria",
        [
            "24/7 log collection from all sources",
            "< 5 minute mean time to detect",
            "95%+ parsing accuracy",
            "Zero critical vulnerabilities",
            "Full team certification",
        ],
        left=0.55,
        top=4.0,
        width=5.8,
        card_color=COLORS["success_green"],
    )

    # Right side: Key metrics boxes
    add_metric_box(
        slide, "24", "Log Sources", left=7.0, top=1.5, color=COLORS["primary_blue"]
    )
    add_metric_box(
        slide,
        "156",
        "Detection Rules",
        left=10.0,
        top=1.5,
        color=COLORS["secondary_blue"],
    )
    add_metric_box(
        slide, "12", "Week Timeline", left=7.0, top=3.0, color=COLORS["tertiary_blue"]
    )
    add_metric_box(
        slide, "8", "Team Members", left=10.0, top=3.0, color=RGBColor(0x00, 0x77, 0xB6)
    )

    # Right side: Compact timeline
    timeline_svg = create_timeline_compact(
        phases=[
            {"name": "Discovery", "status": "complete"},
            {"name": "Design", "status": "in_progress"},
            {"name": "Build", "status": "pending"},
            {"name": "Test", "status": "pending"},
            {"name": "Deploy", "status": "pending"},
        ],
        width=540,
        height=110,
    )
    add_svg_to_slide(
        slide,
        timeline_svg,
        left=6.95,
        top=4.8,
        width=5.9,
        diagram_name="compact_timeline",
    )

    # NEW: Deliverables Detail - Split Layout 60/40
    slide = create_split_slide_60_40(
        prs, "Deliverables Breakdown", "Detailed scope components"
    )

    # Left side (60%): Detailed deliverables table
    add_detail_table(
        slide,
        headers=["Deliverable", "Description", "Priority"],
        rows=[
            ["Platform Setup", "Core SIEM installation and configuration", "P1"],
            ["Log Integration", "Connect 24 log sources across 5 categories", "P1"],
            ["Detection Rules", "156 custom rules + threat intelligence", "P1"],
            ["Dashboards", "12 operational dashboards + reports", "P2"],
            ["SOAR Integration", "Automated response playbooks", "P2"],
            ["Training", "Admin + Analyst certification courses", "P2"],
            ["Documentation", "Runbooks, SOPs, architecture docs", "P3"],
        ],
        left=0.55,
        top=1.5,
        width=7.1,
        font_size=10,
    )

    # Right side (40%): Category breakdown chart
    bar_chart_svg = create_vertical_bar_chart(
        data=[
            {"label": "Network", "value": 92, "max_value": 100},
            {"label": "Endpoint", "value": 88, "max_value": 100},
            {"label": "Cloud", "value": 75, "max_value": 100},
            {"label": "Identity", "value": 95, "max_value": 100},
            {"label": "Apps", "value": 70, "max_value": 100},
        ],
        width=380,
        height=250,
    )
    add_svg_to_slide(
        slide,
        bar_chart_svg,
        left=8.2,
        top=1.5,
        width=4.5,
        diagram_name="coverage_chart",
    )

    # Right side: Status indicators
    add_status_indicator(slide, "complete", "Requirements Gathered", left=8.3, top=5.0)
    add_status_indicator(slide, "complete", "Architecture Approved", left=8.3, top=5.35)
    add_status_indicator(
        slide, "in_progress", "Infrastructure Ready", left=8.3, top=5.7
    )
    add_status_indicator(slide, "pending", "Log Access Confirmed", left=8.3, top=6.05)

    # Section 3: Technical Architecture
    create_section_divider(prs, 3, "Technical Architecture")

    # Architecture diagram slide - with EDITABLE text boxes
    create_editable_architecture_slide(
        prs, "SIEM Architecture", "Editable diagram - click text to customize"
    )

    # Data flow diagram slide
    slide = create_content_slide(prs, "Data Flow Pipeline")
    flow_svg = create_data_flow_diagram(width=1100, height=340)
    add_svg_to_slide(slide, flow_svg, left=0.4, top=1.6, width=12.5)

    # Log sources
    slide = create_content_slide(prs, "Log Source Summary")
    log_svg = create_pie_chart(size=280, show_legend=True)
    add_svg_to_slide(slide, log_svg, left=0.75, top=1.6, width=3)
    add_table(
        slide,
        [
            ["Category", "Sources", "Est. EPS", "Priority"],
            ["Network Security", "Firewall, IDS/IPS, Proxy", "[X]", "P1"],
            ["Endpoint", "EDR, AV, OS Logs", "[X]", "P1"],
            ["Identity", "AD, Azure AD, PAM", "[X]", "P1"],
            ["Cloud", "AWS/Azure/GCP", "[X]", "P2"],
            ["Application", "Web Apps, Databases", "[X]", "P2"],
        ],
        left=4.5,
        top=1.6,
        width=8,
    )

    # Section 4: Implementation
    create_section_divider(prs, 4, "Implementation Approach")

    # NEW: Magazine-style Type 4 - Implementation Strategy
    create_magazine_slide_type4(
        prs,
        header_left="IMPLEMENTATION APPROACH",
        header_right="METHODOLOGY",
        hero_title="Implementation\nStrategy",
        subtitle="Phased Delivery Model",
        body_text="Our implementation follows a structured phased approach. Each phase builds upon the previous, ensuring stable foundations before adding complexity. This approach minimizes risk while maximizing value delivery at each milestone.",
        side_cards=[
            {
                "title": "Phase 1: Foundation",
                "body": "Infrastructure setup, platform deployment, network configuration, and security hardening. Establishes the core platform capabilities.",
            },
            {
                "title": "Phase 2: Integration",
                "body": "Log source connectivity, parser development, data normalization, and quality validation. Ensures comprehensive data collection.",
            },
            {
                "title": "Phase 3: Operationalize",
                "body": "Detection rules deployment, dashboard creation, playbook development, and team training. Delivers operational capabilities.",
            },
        ],
    )

    # Implementation phases diagram
    slide = create_content_slide(prs, "Implementation Methodology")
    phases_svg = create_implementation_phases(width=1100, height=380)
    add_svg_to_slide(slide, phases_svg, left=0.4, top=1.5, width=12.5)

    # Section 5: Timeline
    create_section_divider(prs, 5, "Timeline & Milestones")

    # Timeline diagram
    slide = create_content_slide(prs, "Project Timeline")
    timeline_svg = create_phase_timeline(width=1100, height=180)
    add_svg_to_slide(slide, timeline_svg, left=0.4, top=1.5, width=12.5)
    add_table(
        slide,
        [
            ["Phase", "Start", "End", "Duration"],
            ["Phase 1: Foundation", "[DATE]", "[DATE]", "[X] weeks"],
            ["Phase 2: Integration", "[DATE]", "[DATE]", "[X] weeks"],
            ["Phase 3: Detection", "[DATE]", "[DATE]", "[X] weeks"],
            ["Phase 4: Transition", "[DATE]", "[DATE]", "[X] weeks"],
        ],
        top=4.0,
    )

    # Gantt chart
    slide = create_content_slide(prs, "Detailed Schedule")
    gantt_svg = create_gantt_chart(width=1100, height=400)
    add_svg_to_slide(slide, gantt_svg, left=0.4, top=1.5, width=12.5)

    # Milestones
    slide = create_content_slide(prs, "Key Milestones")
    milestone_svg = create_milestone_roadmap(width=1100, height=220)
    add_svg_to_slide(slide, milestone_svg, left=0.4, top=1.5, width=12.5)

    # Section 6: Governance
    create_section_divider(prs, 6, "Governance & Communication")

    # Governance model diagram
    slide = create_content_slide(prs, "Governance Model")
    gov_svg = create_governance_model(width=1100, height=430)
    add_svg_to_slide(slide, gov_svg, left=0.4, top=1.4, width=12.5)

    # RACI Matrix
    slide = create_content_slide(prs, "RACI Matrix")
    raci_svg = create_raci_matrix(width=1000, height=380)
    add_svg_to_slide(slide, raci_svg, left=0.9, top=1.5, width=11.5)

    # Section 7: Risks
    create_section_divider(prs, 7, "Risks & Mitigation")

    # NEW: Risk Assessment - Split Layout 50/50
    slide = create_split_slide_50_50(
        prs, "Risk Assessment", "Identified risks and mitigation strategies"
    )

    # Left side: Risk matrix diagram
    risk_matrix_svg = create_risk_matrix(width=400, height=340)
    add_svg_to_slide(
        slide,
        risk_matrix_svg,
        left=0.55,
        top=1.5,
        width=5.5,
        diagram_name="risk_matrix",
    )

    # Right side: Risk details cards
    add_info_card(
        slide,
        "High Priority Risks",
        [
            "R1: Infrastructure delays - Early vendor engagement",
            "R4: Scope creep - Strict change control process",
        ],
        left=6.95,
        top=1.5,
        width=5.8,
        card_color=COLORS["error_red"],
    )

    add_info_card(
        slide,
        "Medium Priority Risks",
        [
            "R2: Resource unavailability - Cross-trained backups",
            "R3: Log source access - Pre-validation checklist",
            "R5: Integration complexity - Technical POC phase",
        ],
        left=6.95,
        top=3.3,
        width=5.8,
        card_color=COLORS["warning_yellow"],
    )

    add_info_card(
        slide,
        "Risk Monitoring",
        [
            "Weekly risk review meetings",
            "Automated health checks",
            "Escalation within 24 hours",
        ],
        left=6.95,
        top=5.4,
        width=5.8,
        card_color=COLORS["primary_blue"],
    )

    slide = create_content_slide(prs, "Risk Register")
    add_table(
        slide,
        [
            ["#", "Risk", "Probability", "Impact", "Mitigation"],
            ["R1", "Infrastructure delays", "Medium", "High", "Early engagement"],
            ["R2", "Resource unavailability", "Medium", "Medium", "Backup resources"],
            ["R3", "Log source access issues", "High", "Medium", "Pre-validation"],
            ["R4", "Scope creep", "Medium", "High", "Change control"],
            ["R5", "Integration complexity", "Medium", "Medium", "Technical POC"],
        ],
    )

    # Escalation pyramid
    slide = create_content_slide(prs, "Escalation Path")
    esc_svg = create_escalation_pyramid(width=500, height=350)
    add_svg_to_slide(slide, esc_svg, left=4.2, top=1.5, width=5)

    # Section 8: Next Steps
    create_section_divider(prs, 8, "Next Steps & Q&A")

    # NEW: Magazine-style Type 2 - Key Actions
    create_magazine_slide_type2(
        prs,
        header_left="NEXT STEPS",
        header_right="ACTION ITEMS",
        hero_title="Key Actions\nMoving Forward",
        cards=[
            {
                "title": "Infrastructure Setup",
                "body": "Finalize hardware requirements and submit procurement requests. Ensure network connectivity and firewall rules are in place.",
                "colored": False,
            },
            {
                "title": "Access & Credentials",
                "body": "Prepare log source access credentials. Set up service accounts and API keys for automated collection.",
                "colored": False,
            },
            {
                "title": "Team Coordination",
                "body": "Schedule Phase 1 kickoff meeting. Assign team members and confirm availability for training sessions.",
                "colored": True,
            },
            {
                "title": "Documentation",
                "body": "Distribute meeting minutes and project charter. Review and sign off on scope document.",
                "colored": True,
            },
        ],
    )

    slide = create_content_slide(prs, "Immediate Next Steps")
    add_table(
        slide,
        [
            ["#", "Action", "Owner", "Due Date"],
            ["1", "Distribute meeting minutes", "Seekurity PM", "[DATE]"],
            ["2", "Finalize infrastructure requirements", customer_name, "[DATE]"],
            ["3", "Submit firewall change requests", customer_name, "[DATE]"],
            ["4", "Provide log source credentials", customer_name, "[DATE]"],
            ["5", "Schedule Phase 1 kickoff", "Both PMs", "[DATE]"],
        ],
    )

    # Thank you
    create_thank_you_slide(prs)

    # Save
    prs.save(output_path)
    print(f"  Saved: {output_path}")
    return prs


# =============================================================================
# COMPLETION PRESENTATION GENERATOR
# =============================================================================


def generate_completion_presentation(
    output_path, customer_name="[CUSTOMER_NAME]", date="[DATE]"
):
    """Generate the Project Completion presentation with diagrams."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    print("Generating Completion Presentation...")

    # Title
    create_title_slide(
        prs, "SIEM Implementation", "Project Completion Report", customer_name, date
    )

    # Agenda
    slide = create_content_slide(prs, "Agenda")
    add_table(
        slide,
        [
            ["#", "Topic", "Duration"],
            ["1", "Executive Summary", "10 min"],
            ["2", "Project Achievements", "15 min"],
            ["3", "Technical Delivery", "20 min"],
            ["4", "Metrics & Performance", "10 min"],
            ["5", "Lessons Learned", "10 min"],
            ["6", "Support & Maintenance", "10 min"],
            ["7", "Recommendations", "10 min"],
            ["8", "Sign-off & Closure", "5 min"],
        ],
    )

    # Section 1: Executive Summary
    create_section_divider(prs, 1, "Executive Summary")

    # KPI Dashboard
    slide = create_content_slide(prs, "Key Outcomes", "Project Status: COMPLETED")
    kpi_svg = create_kpi_dashboard(width=1100, height=180)
    add_svg_to_slide(slide, kpi_svg, left=0.4, top=1.6, width=12.5)
    add_table(
        slide,
        [
            ["Metric", "Target", "Achieved", "Status"],
            ["Log Sources", "[X]", "[X]", "100%"],
            ["Use Cases", "[X]", "[X]", "100%"],
            ["Training", "[X] sessions", "[X] sessions", "100%"],
            ["On-Time", "[DATE]", "[DATE]", "On Schedule"],
        ],
        top=4.5,
    )

    # Value delivered - comparison chart
    slide = create_content_slide(prs, "Value Delivered")
    comparison_svg = create_comparison_chart(width=500, height=300)
    add_svg_to_slide(slide, comparison_svg, left=0.75, top=1.6, width=5.5)
    status_svg = create_status_summary(width=450, height=280)
    add_svg_to_slide(slide, status_svg, left=7, top=1.6, width=5.5)

    # NEW: Executive Dashboard - Dashboard Layout with metrics
    slide = create_dashboard_slide(
        prs, "Executive Dashboard", "Project completion at a glance"
    )

    # Top row: Key metric boxes
    add_metric_box(
        slide,
        "100%",
        "Scope Delivered",
        left=0.5,
        top=1.3,
        width=3.0,
        color=COLORS["success_green"],
    )
    add_metric_box(
        slide,
        "24/24",
        "Log Sources",
        left=3.7,
        top=1.3,
        width=3.0,
        color=COLORS["primary_blue"],
    )
    add_metric_box(
        slide,
        "156",
        "Detection Rules",
        left=6.9,
        top=1.3,
        width=3.0,
        color=COLORS["secondary_blue"],
    )
    add_metric_box(
        slide,
        "0",
        "Open Issues",
        left=10.1,
        top=1.3,
        width=2.8,
        color=COLORS["success_green"],
    )

    # Bottom left: Donut charts row
    donut1_svg = create_donut_chart(100, 100, "Scope", width=160, height=160)
    add_svg_to_slide(
        slide, donut1_svg, left=0.6, top=2.8, width=2.0, diagram_name="donut_scope"
    )

    donut2_svg = create_donut_chart(98, 100, "SLA Met", width=160, height=160)
    add_svg_to_slide(
        slide, donut2_svg, left=2.9, top=2.8, width=2.0, diagram_name="donut_sla"
    )

    donut3_svg = create_donut_chart(95, 100, "Uptime", width=160, height=160)
    add_svg_to_slide(
        slide, donut3_svg, left=5.2, top=2.8, width=2.0, diagram_name="donut_uptime"
    )

    # Bottom right: Stats list
    stats_svg = create_stat_list(
        stats=[
            {"label": "Total Events Processed", "value": "1.2B"},
            {"label": "Avg Detection Time", "value": "4.2 min"},
            {"label": "False Positive Rate", "value": "< 5%"},
            {"label": "Team Certified", "value": "12 users"},
            {"label": "Documentation Pages", "value": "145"},
        ],
        width=450,
        height=210,
    )
    add_svg_to_slide(
        slide, stats_svg, left=7.6, top=2.7, width=5.2, diagram_name="stats_list"
    )

    # Timeline at bottom
    timeline_svg = create_timeline_compact(
        phases=[
            {"name": "Discovery", "status": "complete"},
            {"name": "Design", "status": "complete"},
            {"name": "Build", "status": "complete"},
            {"name": "Test", "status": "complete"},
            {"name": "Deploy", "status": "complete"},
        ],
        width=1100,
        height=100,
    )
    add_svg_to_slide(
        slide,
        timeline_svg,
        left=0.5,
        top=5.2,
        width=12.3,
        diagram_name="project_timeline",
    )

    # Section 2: Achievements
    create_section_divider(prs, 2, "Project Achievements")

    # NEW: Project Phases Summary - Three Column Layout
    slide = create_three_column_slide(
        prs, "Implementation Phases", "What we delivered in each phase"
    )

    # Phase 1 card
    add_phase_card(
        slide,
        1,
        "Foundation",
        "complete",
        [
            "Infrastructure deployed",
            "Platform installed",
            "Network configured",
            "Security hardened",
            "Backup configured",
        ],
        left=0.55,
        top=1.5,
        width=3.75,
    )

    # Phase 2 card
    add_phase_card(
        slide,
        2,
        "Integration",
        "complete",
        [
            "24 log sources connected",
            "Parsers configured",
            "Data normalization",
            "Quality validation",
            "Performance tuning",
        ],
        left=4.65,
        top=1.5,
        width=3.75,
    )

    # Phase 3 card
    add_phase_card(
        slide,
        3,
        "Operationalize",
        "complete",
        [
            "156 detection rules",
            "12 dashboards built",
            "SOC playbooks",
            "Team training",
            "Documentation",
        ],
        left=8.75,
        top=1.5,
        width=3.75,
    )

    # Milestone completion
    slide = create_content_slide(prs, "Milestone Completion")
    milestone_svg = create_milestone_roadmap(
        milestones=[
            {"name": "Kickoff", "date": "Week 1", "status": "complete"},
            {"name": "Platform\nDeployed", "date": "Week 2", "status": "complete"},
            {"name": "Sources\nIntegrated", "date": "Week 6", "status": "complete"},
            {"name": "Use Cases\nDeployed", "date": "Week 9", "status": "complete"},
            {"name": "Training\nComplete", "date": "Week 10", "status": "complete"},
            {"name": "Go-Live", "date": "Week 11", "status": "complete"},
        ],
        width=1100,
        height=220,
    )
    add_svg_to_slide(slide, milestone_svg, left=0.4, top=1.5, width=12.5)

    add_table(
        slide,
        [
            ["#", "Milestone", "Target", "Actual", "Status"],
            ["M1", "Infrastructure Ready", "[DATE]", "[DATE]", "PASS"],
            ["M2", "Platform Deployed", "[DATE]", "[DATE]", "PASS"],
            ["M3", "Sources Integrated", "[DATE]", "[DATE]", "PASS"],
            ["M4", "Use Cases Deployed", "[DATE]", "[DATE]", "PASS"],
            ["M5", "Go-Live", "[DATE]", "[DATE]", "PASS"],
        ],
        top=4.2,
        font_size=11,
    )

    # Section 3: Technical Delivery
    create_section_divider(prs, 3, "Technical Delivery")

    # Architecture deployed - with EDITABLE text boxes
    create_editable_architecture_slide(
        prs, "Architecture Deployed", "Editable diagram - click text to customize"
    )

    # Data flow
    slide = create_content_slide(prs, "Data Pipeline Performance")
    flow_svg = create_data_flow_diagram(width=1100, height=340)
    add_svg_to_slide(slide, flow_svg, left=0.4, top=1.6, width=12.5)

    # Log sources
    slide = create_content_slide(prs, "Log Source Integration")
    pie_svg = create_pie_chart(size=300)
    add_svg_to_slide(slide, pie_svg, left=0.5, top=1.5, width=3.5)
    add_table(
        slide,
        [
            ["Category", "Planned", "Integrated", "EPS"],
            ["Network Security", "[X]", "[X]", "[X]"],
            ["Endpoint", "[X]", "[X]", "[X]"],
            ["Identity", "[X]", "[X]", "[X]"],
            ["Cloud", "[X]", "[X]", "[X]"],
            ["Applications", "[X]", "[X]", "[X]"],
            ["Total", "[X]", "[X]", "[X]"],
        ],
        left=4.5,
        top=1.6,
        width=8.3,
    )

    # Section 4: Metrics
    create_section_divider(prs, 4, "Metrics & Performance")

    # KPI Dashboard full
    slide = create_content_slide(prs, "Platform Performance")
    kpi_svg = create_kpi_dashboard(
        kpis=[
            {
                "label": "Uptime",
                "value": "99.9",
                "target": "99.5",
                "unit": "%",
                "status": "success",
            },
            {
                "label": "Search Time",
                "value": "2.1",
                "target": "5",
                "unit": "s",
                "status": "success",
            },
            {
                "label": "Index Latency",
                "value": "18",
                "target": "30",
                "unit": "s",
                "status": "success",
            },
            {
                "label": "Storage",
                "value": "45",
                "target": "80",
                "unit": "%",
                "status": "success",
            },
            {
                "label": "CPU (avg)",
                "value": "42",
                "target": "70",
                "unit": "%",
                "status": "success",
            },
            {
                "label": "Memory",
                "value": "68",
                "target": "85",
                "unit": "%",
                "status": "success",
            },
        ],
        width=1100,
        height=180,
    )
    add_svg_to_slide(slide, kpi_svg, left=0.4, top=1.5, width=12.5)

    # Progress rings
    slide = create_content_slide(prs, "Detection Effectiveness")
    for i, (val, label) in enumerate(
        [(92, "Detection\nCoverage"), (88, "Alert\nFidelity"), (95, "Parse\nRate")]
    ):
        ring_svg = create_progress_ring(val, 100, 150, label.replace("\n", " "))
        add_svg_to_slide(slide, ring_svg, left=1.5 + i * 4, top=2, width=2.5)

    # Comparison
    slide = create_content_slide(prs, "Before vs After")
    comparison_svg = create_comparison_chart(width=800, height=350)
    add_svg_to_slide(slide, comparison_svg, left=2.5, top=1.5, width=8)

    # Section 5: Lessons Learned
    create_section_divider(prs, 5, "Lessons Learned")

    # NEW: Magazine-style Type 3 - Key Learnings
    create_magazine_slide_type3(
        prs,
        header_left="LESSONS LEARNED",
        header_right="PROJECT RETROSPECTIVE",
        quote_text="A successful SIEM implementation requires strong collaboration between security, IT, and business teams. Clear communication and early stakeholder engagement are critical success factors.",
        section_title="Key Takeaways for\nFuture Projects",
        bullet_points=[
            "Start credential provisioning early in the project lifecycle",
            "Regular status meetings maintain alignment across teams",
            "Pre-validation checklists reduce integration issues by 60%",
            "Buffer time in FVT schedule prevents last-minute pressure",
            "Recording training sessions enables on-demand learning",
            "Detailed change requests accelerate approval processes",
        ],
    )

    slide = create_content_slide(prs, "What Went Well")
    add_table(
        slide,
        [
            ["Category", "Observation"],
            ["Planning", "Thorough scoping minimized scope creep"],
            ["Communication", "Regular status meetings kept all parties aligned"],
            ["Technical", "Pre-validation checklist reduced integration issues"],
            ["Collaboration", "Strong customer engagement accelerated decisions"],
            ["Quality", "Iterative testing caught issues early"],
        ],
    )

    slide = create_content_slide(prs, "Areas for Improvement")
    add_table(
        slide,
        [
            ["Category", "Observation", "Recommendation"],
            ["Access", "Credential provisioning delays", "Earlier IAM engagement"],
            ["Testing", "FVT window was tight", "Allocate buffer for FVT"],
            ["Training", "Some missed sessions", "Record for on-demand"],
            ["Change Control", "Some requests lacked detail", "Improve CR template"],
        ],
    )

    # Section 6: Support
    create_section_divider(prs, 6, "Support & Maintenance")

    # Escalation
    slide = create_content_slide(prs, "Support Escalation Path")
    esc_svg = create_escalation_pyramid(width=500, height=350)
    add_svg_to_slide(slide, esc_svg, left=4, top=1.5, width=5.5)

    slide = create_content_slide(prs, "Support SLAs")
    add_table(
        slide,
        [
            ["Severity", "Description", "Response", "Resolution"],
            ["P1 - Critical", "Platform down", "1 hour", "4 hours"],
            ["P2 - High", "Major feature impacted", "4 hours", "8 hours"],
            ["P3 - Medium", "Minor feature impacted", "8 hours", "24 hours"],
            ["P4 - Low", "General questions", "24 hours", "72 hours"],
        ],
    )

    # Section 7: Recommendations
    create_section_divider(prs, 7, "Recommendations")

    slide = create_content_slide(prs, "Roadmap Recommendations")
    roadmap_svg = create_phase_timeline(
        phases=[
            {
                "name": "Optimize",
                "duration": "Q1",
                "color": "#009d9a",
                "status": "active",
            },
            {
                "name": "Expand",
                "duration": "Q2",
                "color": "#0f62fe",
                "status": "pending",
            },
            {
                "name": "Automate",
                "duration": "Q3",
                "color": "#8a3ffc",
                "status": "pending",
            },
            {
                "name": "Mature",
                "duration": "Q4",
                "color": "#0353e9",
                "status": "pending",
            },
        ],
        width=1100,
        height=180,
    )
    add_svg_to_slide(slide, roadmap_svg, left=0.4, top=1.5, width=12.5)
    add_table(
        slide,
        [
            ["#", "Recommendation", "Priority", "Timeframe"],
            ["1", "Complete use case tuning", "High", "30 days"],
            ["2", "Implement SOAR automation", "High", "Q1"],
            ["3", "Expand log source coverage", "Medium", "Q2"],
            ["4", "Deploy behavior analytics", "Medium", "Q3"],
            ["5", "Establish threat hunting program", "Medium", "Q4"],
        ],
        top=4.2,
    )

    # Section 8: Sign-off
    create_section_divider(prs, 8, "Sign-off & Closure")

    slide = create_content_slide(prs, "Acceptance Criteria Verification")
    add_table(
        slide,
        [
            ["#", "Criteria", "Target", "Actual", "Status"],
            ["1", "Log sources integrated", "[X]", "[X]", "PASS"],
            ["2", "Parse success rate", ">95%", "[X]%", "PASS"],
            ["3", "Use cases deployed", "[X]", "[X]", "PASS"],
            ["4", "False positive rate", "<20%", "[X]%", "PASS"],
            ["5", "Platform uptime", ">99.5%", "[X]%", "PASS"],
            ["6", "Training completed", "[X] staff", "[X] staff", "PASS"],
            ["7", "Documentation delivered", "[X] docs", "[X] docs", "PASS"],
        ],
    )

    # Thank you
    create_thank_you_slide(prs)

    # Save
    prs.save(output_path)
    print(f"  Saved: {output_path}")
    return prs


# =============================================================================
# MAIN
# =============================================================================

if __name__ == "__main__":
    script_dir = os.path.dirname(os.path.abspath(__file__))

    print("=" * 60)
    print("Seekurity SIEM Presentation Generator v2")
    print("Enhanced with SVG Diagrams (using svglib)")
    print("=" * 60)
    print()

    kickoff_path = os.path.join(
        script_dir, "SeekuritySIEM_Project_Kickoff_CUSTOMER_NAME.pptx"
    )
    completion_path = os.path.join(
        script_dir, "SeekuritySIEM_Project_Completion_CUSTOMER_NAME.pptx"
    )

    try:
        generate_kickoff_presentation(kickoff_path)
        print()

        generate_completion_presentation(completion_path)
    finally:
        # Clean up Selenium driver
        _cleanup_selenium()

    print()
    print("=" * 60)
    print("Generation complete!")
    print("=" * 60)
